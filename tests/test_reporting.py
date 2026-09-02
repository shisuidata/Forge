from __future__ import annotations

import json
from copy import deepcopy
from datetime import datetime, timedelta, timezone
from pathlib import Path

import pytest

from forge.reporting import ReportStore


def _source() -> dict:
    return {
        "report_id": "rp_demo001",
        "task_run_id": "tr_demo001",
        "org_id": "org_demo",
        "team_id": "team_demo",
        "user_id": "user_demo",
        "revision": 1,
        "bundle_hash": "sha256:" + "a" * 64,
        "title": "地区销售额分析",
        "business_report": {
            "title": "地区销售额分析",
            "executive_summary": "华东销售额最高，但当前结论仅覆盖已审核的查询时间范围。",
            "key_findings": [{
                "statement": "华东销售额最高。", "interpretation": "当前样本中贡献领先。",
                "evidence_refs": ["qr_demo#row:1"], "confidence": "high",
            }],
            "recommendations": [{
                "action": "继续关注华东", "rationale": "贡献最高", "priority": "high",
            }],
            "limitations": ["仅覆盖当前查询时间范围"],
            "next_steps": ["补充退款后的净销售额口径"],
        },
        "analysis": {
            "method_summary": {
                "objective": "比较各地区销售额",
                "dimensions": ["region"],
                "comparison_baseline": "地区横向对比",
                "approach_steps": ["按地区汇总", "比较销售额"],
            },
            "summary": "华东销售额最高。",
            "findings": [{"statement": "华东销售额最高。"}],
            "limitations": ["仅覆盖当前查询时间范围"],
        },
        "query_result": {
            "columns": ["region", "sales"],
            "rows": [["华东", 120], ["华南", 80]],
        },
        "charts": [{
            "chart_type": "bar", "title": "各地区销售额", "dimension": "region",
            "measures": ["sales"], "alt_text": "各地区销售额柱状图",
        }],
        "technical_report": {
            "title": "地区销售额分析 · 技术报告",
            "sql": "SELECT region, SUM(amount) AS sales FROM orders GROUP BY region",
            "approval": {"approved": True},
            "execution": {"row_count": 2},
            "lineage": {"registry_version": "v1"},
            "decision_log": [{"stage": "analysis", "decision": "按地区比较", "rationale": "用户问题要求地区维度"}],
        },
    }


def test_report_store_list_is_scope_aware_and_cursor_stable(tmp_path: Path):
    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    first = _source()
    first["report_id"] = "rp_list001"
    first["task_run_id"] = "tr_list001"
    second = deepcopy(first)
    second["report_id"] = "rp_list002"
    second["task_run_id"] = "tr_list002"
    other = deepcopy(first)
    other["report_id"] = "rp_other001"
    other["task_run_id"] = "tr_other001"
    other["org_id"] = "org_other"
    store.create(first)
    store.create(second)
    store.create(other)

    page = store.list(
        org_id="org_demo", team_id="team_demo", user_id="user_demo", limit=1
    )
    assert len(page) == 1
    assert page[0]["report_id"] in {"rp_list001", "rp_list002"}
    next_page = store.list(
        org_id="org_demo",
        team_id="team_demo",
        user_id="user_demo",
        limit=2,
        before=(page[0]["updated_at"], page[0]["report_id"]),
    )
    assert len(next_page) == 1
    assert {page[0]["report_id"], next_page[0]["report_id"]} == {
        "rp_list001", "rp_list002"
    }
    assert store.list(
        org_id="org_demo", team_id="team_demo", user_id="user_other", limit=10
    ) == []
    with pytest.raises(ValueError, match="unsupported report status"):
        store.list(
            org_id="org_demo", team_id="team_demo", user_id="user_demo", status="ready"
        )


def test_report_store_builds_immutable_html_and_pptx(tmp_path: Path):
    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    created = store.create(_source())
    assert created["status"] == "publishing"

    published = store.build("rp_demo001")

    assert published["status"] == "published"
    assert published["pptx_status"] == "ready"
    html = store.file("rp_demo001", "index.html").read_text()
    technical = store.file("rp_demo001", "technical.html").read_text()
    assert "华东销售额最高" in html
    assert "EXECUTIVE SUMMARY" in html
    assert 'class="finding-card"' in html
    assert "高置信" in html
    assert "高优先级" in html
    assert "限制与风险" in html
    assert "补充退款后的净销售额口径" in html
    assert "@media print" in html
    assert "<svg" in html
    assert "SELECT region" in technical
    assert "hidden chain-of-thought" in technical
    assert (tmp_path / "artifacts").stat().st_mode & 0o777 == 0o700
    assert (tmp_path / "reports.db").stat().st_mode & 0o077 == 0
    from pptx import Presentation
    deck = Presentation(store.file("rp_demo001", "report.pptx"))
    slide_text = "\n".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    for expected in ("执行摘要", "分析思路", "关键发现", "建议行动", "限制与风险", "下一步"):
        assert expected in slide_text
    assert "补充退款后的净销售额口径" in slide_text
    assert len(deck.slides) >= 7
    import sqlite3
    with sqlite3.connect(tmp_path / "reports.db") as conn:
        attempts = conn.execute("SELECT stage, status FROM report_attempts ORDER BY started_at").fetchall()
    assert ("html", "succeeded") in attempts
    assert ("pptx", "succeeded") in attempts


def test_pdf_export_disables_browser_headers_and_internal_path_footer(monkeypatch, tmp_path: Path):
    from types import SimpleNamespace

    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    report_dir = tmp_path / "artifacts" / "rp_demo001" / "v1"
    report_dir.mkdir(parents=True)
    (report_dir / "index.html").write_text("<title>Internal title</title><p>Report</p>")
    captured: list[str] = []

    monkeypatch.setattr("forge.reporting.shutil.which", lambda name: "/usr/bin/google-chrome" if name == "google-chrome" else None)

    def fake_run(command, **_kwargs):
        captured.extend(command)
        output = next(item.split("=", 1)[1] for item in command if item.startswith("--print-to-pdf="))
        Path(output).write_bytes(b"%PDF-1.4\n%%EOF\n")
        return SimpleNamespace(returncode=0)

    monkeypatch.setattr("forge.reporting.subprocess.run", fake_run)
    assert store._build_pdf(report_dir) == "ready"
    assert "--no-pdf-header-footer" in captured
    assert "--print-to-pdf-no-header" not in captured
    assert captured[-1].startswith("file://")


def test_report_projection_suppresses_chart_when_visible_labels_do_not_match_row_grain(tmp_path: Path):
    from pptx import Presentation

    source = _source()
    source["report_id"] = "rp_duplicate001"
    source["task_run_id"] = "tr_duplicate001"
    source["bundle_hash"] = "sha256:" + "d" * 64
    source["query_result"]["rows"] = [["华东", 120], ["华东", 80], ["华南", 60]]
    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    store.create(source)
    assert store.build("rp_duplicate001")["pptx_status"] == "ready"

    report_html = store.file("rp_duplicate001", "index.html").read_text()
    assert "<svg" not in report_html
    assert "各地区销售额</h2>" not in report_html
    deck = Presentation(store.file("rp_duplicate001", "report.pptx"))
    slide_texts = ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in deck.slides]
    assert all("DATA VIEW" not in text for text in slide_texts)


def test_report_projection_escapes_html_and_paginates_long_pptx_content(tmp_path: Path):
    from pptx import Presentation

    source = _source()
    source["report_id"] = "rp_long001"
    source["task_run_id"] = "tr_long001"
    source["bundle_hash"] = "sha256:" + "c" * 64
    source["title"] = source["business_report"]["title"] = "丙" * 220
    source["charts"][0]["title"] = "丁" * 180
    source["business_report"]["executive_summary"] = "<img src=x onerror=alert(1)>" + "乙" * 260
    source["business_report"]["key_findings"][0]["statement"] = "头" + "甲" * 300 + "尾"
    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    published = store.create(source)
    assert published["status"] == "publishing"
    assert store.build("rp_long001")["pptx_status"] == "ready"

    report_html = store.file("rp_long001", "index.html").read_text()
    assert "<img src=x" not in report_html
    assert "&lt;img src=x onerror=alert(1)&gt;" in report_html
    deck = Presentation(store.file("rp_long001", "report.pptx"))
    all_text = "".join(shape.text for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert all_text.count("甲") == 300
    assert all_text.count("丙") >= 220
    assert all_text.count("丁") >= 180
    assert "头" in all_text and "尾" in all_text
    assert max(len(shape.text) for slide in deck.slides for shape in slide.shapes if hasattr(shape, "text")) <= 160
    assert len(deck.slides) >= 11


def test_report_store_is_idempotent_and_rejects_bundle_drift(tmp_path: Path):
    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    first = store.create(_source())
    assert store.create(_source())["report_id"] == first["report_id"]
    changed = _source()
    changed["bundle_hash"] = "sha256:" + "b" * 64
    with pytest.raises(ValueError, match="different bundle"):
        store.create(changed)


def test_report_store_restart_marks_running_attempt_interrupted_without_replay(tmp_path: Path):
    import sqlite3
    db = tmp_path / "reports.db"
    store = ReportStore(str(db), str(tmp_path / "artifacts"))
    store.create(_source())
    with sqlite3.connect(db) as conn:
        conn.execute(
            "INSERT INTO report_attempts VALUES ('rpa_running', 'rp_demo001', 'pdf', 'running', ?, NULL, NULL)",
            (datetime.now(timezone.utc).isoformat(),),
        )
    ReportStore(str(db), str(tmp_path / "artifacts"))
    with sqlite3.connect(db) as conn:
        attempt = conn.execute("SELECT status FROM report_attempts WHERE attempt_id='rpa_running'").fetchone()
        report = conn.execute("SELECT status FROM reports WHERE report_id='rp_demo001'").fetchone()
    assert attempt == ("interrupted",)
    assert report == ("failed",)


def test_report_store_rejects_hidden_reasoning(tmp_path: Path):
    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    source = _source()
    source["analysis"]["summary"] = "<think>private reasoning</think>"
    with pytest.raises(ValueError, match="forbidden reasoning"):
        store.create(source)


@pytest.mark.asyncio
async def test_business_share_is_expiring_revocable_and_cannot_open_technical_report(client, monkeypatch, tmp_path):
    import web.routes.reports as reports
    from config import cfg

    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    store.create(_source())
    store.build("rp_demo001")
    monkeypatch.setattr(reports, "get_report_store", lambda: store)
    share = store.create_share(
        "rp_demo001", (datetime.now(timezone.utc) + timedelta(hours=1)).isoformat()
    )
    old_auth = cfg.AUTH_ENABLED
    cfg.AUTH_ENABLED = True
    try:
        landing = await client.get(f"/reports/share/{share['share_id']}")
        assert landing.status_code == 200
        assert share["token"] not in landing.text
        exchange = await client.post(
            f"/reports/share/{share['share_id']}/exchange", json={"token": share["token"]}
        )
        assert exchange.status_code == 200
        page = await client.get("/reports/rp_demo001")
        assert page.status_code == 200
        technical = await client.get("/reports/rp_demo001/technical")
        assert technical.status_code == 401
        store.revoke_share("rp_demo001", share["share_id"])
        revoked = await client.get("/reports/rp_demo001")
        assert revoked.status_code == 401
    finally:
        cfg.AUTH_ENABLED = old_auth


@pytest.mark.asyncio
async def test_report_api_requires_pi_auth_and_serves_authenticated_outputs(client, monkeypatch, tmp_path):
    import web.routes.reports as reports

    store = ReportStore(str(tmp_path / "reports.db"), str(tmp_path / "artifacts"))
    reports.get_report_store.cache_clear()
    monkeypatch.setattr(reports, "get_report_store", lambda: store)

    unauthorized = await client.post("/api/internal/reports", json=_source())
    assert unauthorized.status_code == 401

    from config import cfg
    old_keys = cfg.PI_SERVICE_API_KEYS
    cfg.PI_SERVICE_API_KEYS = ["pi-secret"]
    try:
        created = await client.post(
            "/api/internal/reports",
            json=_source(),
            headers={"X-Pi-Service-Key": "pi-secret"},
        )
        assert created.status_code == 202
        status = await client.get(
            "/api/internal/reports/rp_demo001",
            headers={"X-Pi-Service-Key": "pi-secret"},
        )
        assert status.json()["report"]["status"] == "published"
        page = await client.get("/reports/rp_demo001")
        assert page.status_code == 200
        assert "地区销售额分析" in page.text
        pptx = await client.get("/reports/rp_demo001/download/pptx")
        assert pptx.status_code == 200
    finally:
        cfg.PI_SERVICE_API_KEYS = old_keys
