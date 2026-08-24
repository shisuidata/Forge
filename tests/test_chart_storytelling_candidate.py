from __future__ import annotations

import importlib.util
import json
from pathlib import Path

import pytest

from agent.contracts import validate_contract

FIXTURES = Path(__file__).parent / "fixtures" / "chart-storytelling"
CANDIDATE_TOOL = Path(__file__).parents[1] / "tools" / "chart_storytelling_candidate.py"
_TOOL_SPEC = importlib.util.spec_from_file_location("chart_storytelling_candidate", CANDIDATE_TOOL)
assert _TOOL_SPEC is not None and _TOOL_SPEC.loader is not None
_TOOL = importlib.util.module_from_spec(_TOOL_SPEC)
_TOOL_SPEC.loader.exec_module(_TOOL)
_html, _load, _pptx = _TOOL._html, _TOOL._load, _TOOL._pptx
_pptx_from_images = _TOOL._pptx_from_images


def _fixture(name: str) -> dict:
    return json.loads((FIXTURES / name).read_text())


def _row_number(ref: str) -> int:
    return int(ref.rsplit(":", 1)[1])


@pytest.mark.parametrize("name", ["category-comparison.json", "time-series.json"])
def test_chart_v2_fixtures_are_schema_valid_and_evidence_bounded(name: str):
    fixture = _fixture(name)
    query = fixture["query_result"]
    columns = set(query["columns"])
    query_prefix = f'{query["query_run_id"]}#row:'
    for chart in fixture["charts"]:
        validate_contract("chart_artifact_v2", chart)
        payload = chart["payload"]
        assert payload["quality_status"]["status"] == "ready"
        assert set(payload["grain"]["key_fields"]).issubset(columns)
        assert payload["grain"]["display_field"] in columns
        refs = set(payload["evidence_refs"])
        assert all(ref.startswith(query_prefix) and 1 <= _row_number(ref) <= query["row_count"] for ref in refs)
        assert all(set(series["evidence_refs"]).issubset(refs) for series in payload["series"])
        assert all(set(annotation["evidence_refs"]).issubset(refs) for annotation in payload["annotations"])
        series_ids = {series["series_id"] for series in payload["series"]}
        assert all(annotation["target"].get("series_id") in series_ids for annotation in payload["annotations"] if annotation["target"].get("series_id"))
        source_fields = columns | {transform.get("output_field") for transform in payload["transforms"] if transform.get("output_field")}
        assert all(series["field"] in source_fields for series in payload["series"])


def test_category_story_recomputes_ranking_pareto_and_top_n_other():
    fixture = _fixture("category-comparison.json")
    rows = fixture["query_result"]["rows"]
    values = [row[1] for row in rows]
    total = sum(values)
    expected = fixture["expected"]
    assert total == expected["total_sales"]
    assert sum(values[:3]) / total == pytest.approx(expected["top_three_share"], abs=0.0001)
    assert sum(values[:6]) / total == pytest.approx(expected["top_six_share"], abs=0.0001)
    assert sum(values[-3:]) / total == pytest.approx(expected["tail_three_share"], abs=0.0001)
    ranking = fixture["charts"][0]["payload"]
    top_n = next(transform for transform in ranking["transforms"] if transform["type"] == "top_n")
    assert top_n == {"type": "top_n", "field": "sales_amount", "limit": 7, "remainder_label": "其他 3 类"}
    assert sum(values[: top_n["limit"]]) + sum(values[top_n["limit"] :]) == total


def test_time_story_recomputes_target_gaps_and_channel_contribution():
    fixture = _fixture("time-series.json")
    rows = fixture["query_result"]["rows"]
    april, june = rows[3], rows[5]
    expected = fixture["expected"]
    assert (april[4] - april[5]) / april[5] == pytest.approx(expected["april_gap_ratio"], abs=0.0001)
    assert (june[4] - june[5]) / june[5] == pytest.approx(expected["june_above_target_ratio"], abs=0.0001)
    growth = june[4] - april[4]
    direct_growth = june[1] - april[1]
    assert growth == expected["april_to_june_growth"]
    assert direct_growth == expected["direct_growth"]
    assert direct_growth / growth == pytest.approx(expected["direct_growth_share"], abs=0.0001)


def test_chart_v2_contract_rejects_markup_and_evidence_without_a_source_row():
    chart = _fixture("category-comparison.json")["charts"][0]
    chart["payload"]["title"] = "<style>unsafe</style>"
    with pytest.raises(Exception):
        validate_contract("chart_artifact_v2", chart)

    chart = _fixture("category-comparison.json")["charts"][0]
    chart["payload"]["annotations"][0]["evidence_refs"] = []
    with pytest.raises(Exception):
        validate_contract("chart_artifact_v2", chart)


def test_candidate_html_is_self_contained_and_preserves_four_decision_views():
    rendered = _html(_load("category-comparison.json"), _load("time-series.json"))
    assert rendered.count('class="chart-card chart-wide"') == 4
    assert rendered.count('class="annotation ') == 4
    assert "头部品类领先，但没有单一绝对赢家" in rendered
    assert "经营六个品类，可覆盖 82.2% 销售额" in rendered
    assert "四月失速，五月完成反转" in rendered
    assert "直营贡献了反转后 50% 的新增量" in rendered
    assert "Y 轴从 650K 起" in rendered
    assert "https://" not in rendered
    assert "<script src=" not in rendered


def test_candidate_pptx_is_static_complete_and_within_slide_bounds(tmp_path: Path):
    from pptx import Presentation

    target = tmp_path / "candidate.pptx"
    _pptx(_load("category-comparison.json"), _load("time-series.json"), target)
    deck = Presentation(target)
    assert len(deck.slides) == 5
    expected = [
        "从图表堆砌",
        "头部品类接近",
        "经营六个品类",
        "四月失速",
        "直营贡献了反转后",
    ]
    for slide, phrase in zip(deck.slides, expected, strict=True):
        text = "\n".join(
            shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        )
        assert phrase in text
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0
            assert shape.left + shape.width <= deck.slide_width
            assert shape.top + shape.height <= deck.slide_height
    final_text = "\n".join(
        shape.text
        for shape in deck.slides[-1].shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert all(label in final_text for label in ("直营", "平台", "门店", "87K", "50%"))


def test_candidate_image_deck_preserves_one_complete_static_slide_per_image(tmp_path: Path):
    from PIL import Image
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    images = []
    for index in range(5):
        path = tmp_path / f"slide-{index}.png"
        Image.new("RGB", (1280, 720), (20 + index, 63, 49)).save(path)
        images.append(path)
    target = tmp_path / "image-deck.pptx"
    _pptx_from_images(images, target)
    deck = Presentation(target)
    assert len(deck.slides) == 5
    for slide in deck.slides:
        assert len(slide.shapes) == 1
        picture = slide.shapes[0]
        assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE
        assert picture.left == 0 and picture.top == 0
        assert picture.width == deck.slide_width and picture.height == deck.slide_height
