#!/usr/bin/env python3
"""Generate the H5 ChartArtifact v2 visual candidates from fixed fixtures.

This is an R0 design gate, not a production report renderer. It consumes only
versioned fixture data and never invokes a model or a database.
"""
from __future__ import annotations

import argparse
import html
import json
import math
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[1]
FIXTURES = ROOT / "tests" / "fixtures" / "chart-storytelling"


def _load(name: str) -> dict[str, Any]:
    return json.loads((FIXTURES / name).read_text())


def _compact(value: float, symbol: str = "¥") -> str:
    if abs(value) >= 1_000_000:
        return f"{symbol}{value / 1_000_000:.2f}M"
    if abs(value) >= 1_000:
        return f"{symbol}{value / 1_000:.0f}K"
    return f"{symbol}{value:,.0f}"


def _evidence_button(label: str, refs: list[str]) -> str:
    safe_label = html.escape(label)
    safe_refs = html.escape(" · ".join(refs), quote=True)
    return f'<button class="evidence-link" type="button" data-evidence="{safe_refs}" data-note="{safe_label}">↗ 查看证据</button>'


def _category_ranking(fixture: dict[str, Any]) -> str:
    rows = fixture["query_result"]["rows"]
    top = rows[:7]
    other = ["其他 3 类", sum(row[1] for row in rows[7:]), sum(row[2] for row in rows[7:])]
    points = top + [other]
    maximum = max(row[1] for row in points)
    width, height, left, right, top_pad, row_gap = 900, 500, 135, 78, 28, 55
    chunks = []
    for index, (label, value, _orders) in enumerate(points):
        y = top_pad + index * row_gap
        bar_width = (width - left - right) * value / maximum
        tone = "rank-top" if index == 0 else "rank-other" if index == len(points) - 1 else "rank-bar"
        evidence = f"qr_category_story#row:{index + 1}" if index < 7 else "qr_category_story#row:8 · row:9 · row:10"
        chunks.append(
            f'<g class="datum" tabindex="0" role="img" aria-label="{html.escape(label)}，销售额 {_compact(value)}，证据 {evidence}">'
            f'<text x="0" y="{y + 17}" class="axis-label">{html.escape(label)}</text>'
            f'<rect x="{left}" y="{y}" width="{bar_width:.1f}" height="28" rx="7" class="{tone}"><title>{html.escape(label)} · {_compact(value)}</title></rect>'
            f'<text x="{left + bar_width + 10:.1f}" y="{y + 19}" class="value-label">{_compact(value)}</text></g>'
        )
    refs = fixture["charts"][0]["payload"]["annotations"][0]["evidence_refs"]
    return (
        '<article class="chart-card chart-wide" data-chart="category-ranking">'
        '<div class="chart-heading"><div><span class="chart-kicker">01 · RANKING</span><h3>头部品类领先，但没有单一绝对赢家</h3>'
        '<p>先看规模排序，再判断资源是否应该只押注第一名。</p></div><span class="unit-pill">销售额 · CNY</span></div>'
        f'<svg class="chart-svg" viewBox="0 0 {width} {height}" aria-label="品类销售额排名">{"".join(chunks)}</svg>'
        '<div class="annotation annotation-lime"><span>关键标注</span><strong>玩具礼品仅比运动户外高 4.4%</strong><p>头部竞争接近，建议采用组合经营而不是单品类押注。</p>'
        f'{_evidence_button("头部差距来自排名前两行", refs)}</div>'
        '<button class="table-toggle" type="button" aria-expanded="false">查看数据表</button>'
        f'{_table(fixture["query_result"], "category-table")}</article>'
    )


def _category_pareto(fixture: dict[str, Any]) -> str:
    rows = fixture["query_result"]["rows"]
    values = [row[1] for row in rows]
    total = sum(values)
    cumulative, running = [], 0
    for value in values:
        running += value
        cumulative.append(running / total)
    width, height, left, right, top, bottom = 900, 460, 62, 44, 38, 82
    inner_w, inner_h = width - left - right, height - top - bottom
    gap = inner_w / len(rows)
    max_value = max(values)
    bars, labels, coords = [], [], []
    for index, row in enumerate(rows):
        x = left + index * gap + 8
        bar_h = inner_h * 0.46 * row[1] / max_value
        bars.append(f'<rect x="{x:.1f}" y="{top + inner_h - bar_h:.1f}" width="{gap - 16:.1f}" height="{bar_h:.1f}" rx="5" class="pareto-bar"><title>{html.escape(row[0])} · {_compact(row[1])}</title></rect>')
        labels.append(f'<text x="{x + (gap - 16) / 2:.1f}" y="{height - 48}" class="mini-label" text-anchor="middle">{html.escape(row[0][:2])}</text>')
        y = top + inner_h * (1 - cumulative[index])
        coords.append((x + (gap - 16) / 2, y))
    path = " ".join(("M" if i == 0 else "L") + f" {x:.1f} {y:.1f}" for i, (x, y) in enumerate(coords))
    dots = "".join(f'<circle cx="{x:.1f}" cy="{y:.1f}" r="{7 if i == 5 else 4}" class="pareto-dot" tabindex="0"><title>累计贡献 {cumulative[i] * 100:.1f}%</title></circle>' for i, (x, y) in enumerate(coords))
    y80 = top + inner_h * 0.2
    percent_ticks = "".join(
        f'<text x="{width-right+5}" y="{top + inner_h * (1-share) + 4:.1f}" class="mini-label">{int(share*100)}%</text>'
        for share in (0.2, 0.4, 0.6, 0.8, 1.0)
    )
    refs = fixture["charts"][1]["payload"]["annotations"][1]["evidence_refs"]
    return (
        '<article class="chart-card chart-wide" data-chart="category-pareto">'
        '<div class="chart-heading"><div><span class="chart-kicker">02 · CONTRIBUTION</span><h3>经营六个品类，可覆盖 82.2% 销售额</h3>'
        '<p>累计贡献回答的是资源覆盖边界，而不是另一张排名图。</p></div><span class="unit-pill">累计贡献 · %</span></div>'
        f'<svg class="chart-svg" viewBox="0 0 {width} {height}" aria-label="品类累计贡献帕累托图">'
        f'<line x1="{left}" y1="{y80:.1f}" x2="{width-right}" y2="{y80:.1f}" class="reference-line"/><text x="{width-right}" y="{y80-10:.1f}" text-anchor="end" class="reference-label">80% 经营覆盖线</text>'
        f'{percent_ticks}{"".join(bars)}<path d="{path}" class="pareto-line"/>{dots}{"".join(labels)}</svg>'
        '<div class="annotation annotation-coral"><span>决策阈值</span><strong>第六个品类跨过 80%</strong><p>头部五类负责规模，第六类决定覆盖线；尾部三类合计仅 10.0%。</p>'
        f'{_evidence_button("累计贡献由前六行复算", refs)}</div></article>'
    )


def _line_chart(fixture: dict[str, Any]) -> str:
    rows = fixture["query_result"]["rows"]
    width, height, left, right, top, bottom = 900, 470, 70, 46, 52, 70
    inner_w, inner_h = width - left - right, height - top - bottom
    values = [row[4] for row in rows] + [row[5] for row in rows]
    low, high = 650000, 930000
    def xy(index: int, value: float) -> tuple[float, float]:
        return left + index * inner_w / (len(rows) - 1), top + inner_h * (high - value) / (high - low)
    actual = [xy(i, row[4]) for i, row in enumerate(rows)]
    target = [xy(i, row[5]) for i, row in enumerate(rows)]
    path = lambda pts: " ".join(("M" if i == 0 else "L") + f" {x:.1f} {y:.1f}" for i, (x, y) in enumerate(pts))
    grid = []
    for value in [700000, 750000, 800000, 850000, 900000]:
        y = xy(0, value)[1]
        grid.append(f'<line x1="{left}" y1="{y:.1f}" x2="{width-right}" y2="{y:.1f}" class="grid-line"/><text x="{left-12}" y="{y+4:.1f}" text-anchor="end" class="mini-label">{value/1000:.0f}K</text>')
    months = "".join(f'<text x="{actual[i][0]:.1f}" y="{height-32}" text-anchor="middle" class="axis-label">{row[0][5:]}月</text>' for i, row in enumerate(rows))
    dots = []
    for series_id, points, index_value in [("actual", actual, 4), ("target", target, 5)]:
        for i, (x, y) in enumerate(points):
            dots.append(f'<circle cx="{x:.1f}" cy="{y:.1f}" r="{7 if (series_id == "actual" and i in (3,5)) else 4}" class="line-dot series-{series_id}" tabindex="0" aria-label="{rows[i][0]} {"实际" if series_id == "actual" else "目标"} {_compact(rows[i][index_value])}"><title>{rows[i][0]} · {_compact(rows[i][index_value])}</title></circle>')
    april_refs = fixture["charts"][0]["payload"]["annotations"][0]["evidence_refs"]
    return (
        '<article class="chart-card chart-wide" data-chart="monthly-trend">'
        '<div class="chart-heading"><div><span class="chart-kicker">03 · TREND</span><h3>四月失速，五月完成反转</h3><p>实际与目标同图，避免只看增长却忽略经营承诺。</p></div><span class="unit-pill">月销售额 · CNY</span></div>'
        '<div class="legend" aria-label="图例"><button type="button" data-series="actual" aria-pressed="true"><i class="legend-actual"></i>实际</button><button type="button" data-series="target" aria-pressed="true"><i class="legend-target"></i>目标</button></div>'
        f'<svg class="chart-svg" viewBox="0 0 {width} {height}" aria-label="月度实际销售额与目标趋势">{"".join(grid)}'
        f'<path d="{path(target)}" class="line-target chart-series series-target"/><path d="{path(actual)}" class="line-actual chart-series series-actual"/>{"".join(dots)}{months}'
        f'<text x="{left}" y="{height-7}" class="axis-note">// Y 轴从 650K 起，用于观察目标偏差</text></svg>'
        '<div class="annotation annotation-coral"><span>异常点</span><strong>四月低于目标 6.4%</strong><p>这是六个月中唯一明显失速点；五月回到目标线上方，六月扩大优势。</p>'
        f'{_evidence_button("四月偏差由当月实际与目标复算", april_refs)}</div>'
        '<button class="table-toggle" type="button" aria-expanded="false">查看数据表</button>'
        f'{_table(fixture["query_result"], "monthly-table")}</article>'
    )


def _stacked_area(fixture: dict[str, Any]) -> str:
    rows = fixture["query_result"]["rows"]
    width, height, left, right, top, bottom = 900, 460, 64, 42, 42, 72
    inner_w, inner_h = width-left-right, height-top-bottom
    high = 930000
    def x(i: int) -> float: return left + i * inner_w / (len(rows)-1)
    def y(v: float) -> float: return top + inner_h * (1-v/high)
    layers = [("retail", 3), ("marketplace", 2), ("direct", 1)]
    cumulative = [0.0]*len(rows); paths=[]
    for series, idx in layers:
        lower=cumulative[:]; upper=[lower[i]+rows[i][idx] for i in range(len(rows))]; cumulative=upper
        top_pts=[(x(i),y(v)) for i,v in enumerate(upper)]; bottom_pts=[(x(i),y(v)) for i,v in reversed(list(enumerate(lower)))]
        d=" ".join(("M" if i==0 else "L")+f" {px:.1f} {py:.1f}" for i,(px,py) in enumerate(top_pts+bottom_pts))+" Z"
        paths.append(f'<path d="{d}" class="area-series area-{series}" data-area-series="{series}"><title>{series}</title></path>')
    months="".join(f'<text x="{x(i):.1f}" y="{height-30}" text-anchor="middle" class="axis-label">{row[0][5:]}月</text>' for i,row in enumerate(rows))
    final_total = rows[-1][4]
    final_market_retail = rows[-1][2] + rows[-1][3]
    final_retail = rows[-1][3]
    label_x = x(len(rows)-1) - 4
    direct_label_y = y(final_total) + 18
    market_label_y = y(final_market_retail) + 18
    retail_label_y = y(final_retail) + 18
    april_x, june_x = x(3), x(5)
    april_direct_top = y(rows[3][4])
    june_direct_top = y(rows[5][4])
    refs=fixture["charts"][1]["payload"]["annotations"][0]["evidence_refs"]
    return (
        '<article class="chart-card chart-wide" data-chart="monthly-mix">'
        '<div class="chart-heading"><div><span class="chart-kicker">04 · COMPOSITION</span><h3>直营贡献了反转后 50% 的新增量</h3><p>结构图解释“增长从哪里来”，不重复趋势图的时间结论。</p></div><span class="unit-pill">渠道结构 · CNY</span></div>'
        '<div class="legend" aria-label="渠道图例"><button type="button" data-area="direct" aria-pressed="true"><i class="legend-direct"></i>直营</button><button type="button" data-area="marketplace" aria-pressed="true"><i class="legend-market"></i>平台</button><button type="button" data-area="retail" aria-pressed="true"><i class="legend-retail"></i>门店</button></div>'
        f'<svg class="chart-svg" viewBox="0 0 {width} {height}" aria-label="月度渠道销售额结构"><defs><marker id="growth-arrow" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="6" markerHeight="6" orient="auto-start-reverse"><path d="M 0 0 L 10 5 L 0 10 z" class="arrow-head"/></marker></defs>{"".join(paths)}{months}'
        f'<text x="{label_x:.1f}" y="{direct_label_y:.1f}" text-anchor="end" class="area-label area-label-light">直营</text><text x="{label_x:.1f}" y="{market_label_y:.1f}" text-anchor="end" class="area-label">平台</text><text x="{label_x:.1f}" y="{retail_label_y:.1f}" text-anchor="end" class="area-label">门店</text>'
        f'<line x1="{april_x:.1f}" y1="{april_direct_top-18:.1f}" x2="{june_x-12:.1f}" y2="{june_direct_top-18:.1f}" class="growth-arrow" marker-end="url(#growth-arrow)"/><text x="{(april_x+june_x)/2:.1f}" y="{min(april_direct_top,june_direct_top)-28:.1f}" text-anchor="middle" class="growth-label">直营 +87K</text></svg>'
        '<div class="annotation annotation-lime"><span>贡献拆解</span><strong>四月至六月新增 174K，直营贡献 87K</strong><p>恢复不是全渠道平均发生；直营是首要验证与加码对象。</p>'
        f'{_evidence_button("增量贡献由四月与六月两行复算", refs)}</div></article>'
    )


def _table(query: dict[str, Any], table_id: str) -> str:
    head="".join(f"<th>{html.escape(str(c))}</th>" for c in query["columns"])
    body="".join("<tr>"+"".join(f"<td>{html.escape(str(v))}</td>" for v in row)+"</tr>" for row in query["rows"])
    return f'<div class="data-table" id="{table_id}" hidden><table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table></div>'


def _html(category: dict[str, Any], monthly: dict[str, Any]) -> str:
    category_total=category["expected"]["total_sales"]
    june=monthly["query_result"]["rows"][-1]
    template="""<!doctype html><html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>Forge · Evidence Story R0</title><style>
:root{--ink:#13231e;--paper:#eeeadd;--card:#fffdf7;--moss:#123f31;--green:#2f745d;--mint:#8ad3ad;--lime:#c8ef72;--coral:#e27a5f;--amber:#d39a35;--slate:#718079;--line:#d7ddd7;--soft:#edf2ed}*{box-sizing:border-box}html{scroll-behavior:smooth}body{margin:0;background:radial-gradient(circle at 80% 0,#c8ef7230,transparent 28rem),linear-gradient(135deg,#123f3107 25%,transparent 25%) 0 0/22px 22px,var(--paper);color:var(--ink);font:15px/1.65 "Avenir Next","PingFang SC",sans-serif}.topbar{position:sticky;top:0;z-index:10;display:flex;justify-content:space-between;align-items:center;padding:12px max(24px,calc((100vw - 1240px)/2));background:#eeeadddf;backdrop-filter:blur(15px);border-bottom:1px solid #13231e1a}.brand{font:800 12px/1 "Avenir Next",sans-serif;letter-spacing:.18em}.brand b{color:var(--green)}nav{display:flex;gap:8px}nav a,.ghost-btn{border:1px solid #13231e24;border-radius:999px;background:#fffdf7;padding:7px 12px;color:var(--ink);text-decoration:none;font-size:12px;font-weight:700}.hero{max-width:1240px;margin:32px auto 18px;border-radius:30px;background:var(--moss);color:#f7fbf8;padding:62px 68px;position:relative;overflow:hidden}.hero:after{content:"";position:absolute;right:-80px;top:-180px;width:450px;height:450px;border:1px solid #c8ef7266;border-radius:50%;box-shadow:0 0 0 64px #c8ef7209,0 0 0 128px #c8ef7207}.eyebrow,.scenario-kicker,.chart-kicker{font-size:11px;font-weight:900;letter-spacing:.18em;color:var(--mint)}.hero h1{max-width:900px;margin:14px 0 18px;font:700 clamp(42px,6vw,76px)/1.02 "Iowan Old Style","Songti SC",serif;letter-spacing:-.045em}.hero p{max-width:740px;margin:0;color:#d7e8e0;font-size:19px}.hero .rule{width:78px;height:6px;background:var(--lime);margin:28px 0}.hero-meta{display:flex;gap:18px;flex-wrap:wrap;margin-top:30px}.hero-meta span{border:1px solid #ffffff2b;border-radius:999px;padding:8px 13px;font-size:12px}.report{max-width:1240px;margin:auto;padding:0 0 100px}.scenario{margin-top:24px;border:1px solid #13231e16;border-radius:26px;background:#f8f6efcc;padding:34px}.scenario-head{display:grid;grid-template-columns:1fr auto;gap:30px;align-items:end}.scenario h2{max-width:800px;margin:6px 0 0;font:700 clamp(32px,4vw,50px)/1.08 "Iowan Old Style","Songti SC",serif;letter-spacing:-.03em}.quality{display:inline-flex;gap:8px;align-items:center;border-radius:999px;background:#e8f7ed;color:#17613f;padding:8px 12px;font-size:11px;font-weight:800}.quality:before{content:"";width:7px;height:7px;border-radius:50%;background:#19a66d}.metrics{display:grid;grid-template-columns:repeat(3,1fr);gap:12px;margin:26px 0}.metric{border:1px solid var(--line);border-radius:18px;background:var(--card);padding:18px 20px}.metric span{display:block;color:var(--slate);font-size:11px;font-weight:800;letter-spacing:.08em}.metric strong{display:block;margin-top:4px;font:700 28px/1.1 "Iowan Old Style",serif}.metric p{margin:7px 0 0;color:var(--slate);font-size:12px}.chart-grid{display:grid;grid-template-columns:1fr 1fr;gap:16px}.chart-card{position:relative;border:1px solid var(--line);border-radius:22px;background:var(--card);padding:26px;overflow:hidden}.chart-wide{grid-column:span 1}.chart-heading{display:flex;justify-content:space-between;gap:20px}.chart-heading h3{margin:6px 0 5px;font:700 25px/1.18 "Iowan Old Style","Songti SC",serif}.chart-heading p{margin:0;color:var(--slate);font-size:13px}.unit-pill{align-self:start;white-space:nowrap;border-radius:999px;background:var(--soft);padding:7px 10px;color:#3f554c;font-size:10px;font-weight:800}.chart-svg{display:block;width:100%;height:auto;margin:18px 0 4px;overflow:visible}.axis-label{fill:#4b5c55;font-size:12px}.mini-label{fill:#718079;font-size:10px}.value-label{fill:#253c33;font-size:11px;font-weight:800}.rank-bar{fill:#4f8f78}.rank-top{fill:var(--lime)}.rank-other{fill:#d8a48f}.pareto-bar{fill:#dce9e1}.pareto-line{fill:none;stroke:var(--green);stroke-width:4}.pareto-dot{fill:var(--card);stroke:var(--green);stroke-width:3}.reference-line{stroke:var(--coral);stroke-width:1.5;stroke-dasharray:7 6}.reference-label{fill:var(--coral);font-size:11px;font-weight:800}.grid-line{stroke:#dce2de;stroke-width:1}.line-actual{fill:none;stroke:var(--green);stroke-width:5}.line-target{fill:none;stroke:var(--coral);stroke-width:3;stroke-dasharray:9 8}.line-dot{fill:var(--card);stroke-width:3}.series-actual{stroke:var(--green)}.series-target{stroke:var(--coral)}.area-direct{fill:#2f745d}.area-marketplace{fill:#74b596}.area-retail{fill:#c9dfd1}.area-series{transition:opacity .18s}.area-series.is-muted,.chart-series.is-muted,.line-dot.is-muted{opacity:.12}.axis-note{fill:#7c8983;font-size:9px}.area-label{fill:#173f32;font-size:10px;font-weight:900}.area-label-light{fill:#f4fff8}.growth-arrow{stroke:var(--coral);stroke-width:2;stroke-dasharray:5 4}.arrow-head{fill:var(--coral)}.growth-label{fill:var(--coral);font-size:10px;font-weight:900}.annotation{display:grid;grid-template-columns:auto 1fr auto;gap:4px 15px;align-items:center;border-radius:15px;padding:14px 16px}.annotation>span{grid-row:1/3;align-self:start;border-radius:999px;padding:4px 7px;font-size:9px;font-weight:900;letter-spacing:.08em}.annotation strong{font-size:13px}.annotation p{grid-column:2;margin:0;color:#53665e;font-size:12px}.annotation-lime{background:#eff8d9;border:1px solid #c8e795}.annotation-lime>span{background:var(--lime)}.annotation-coral{background:#fff0eb;border:1px solid #efc4b8}.annotation-coral>span{background:#f3b7a7}.evidence-link{grid-column:3;grid-row:1/3;border:0;background:transparent;color:#1f6b52;text-decoration:underline;text-underline-offset:3px;font-size:11px;font-weight:800;cursor:pointer}.legend{display:flex;gap:8px;margin-top:15px}.legend button{border:1px solid var(--line);border-radius:999px;background:white;padding:6px 9px;font-size:11px;font-weight:700;cursor:pointer}.legend button[aria-pressed="false"]{opacity:.42}.legend i{display:inline-block;width:8px;height:8px;border-radius:50%;margin-right:5px}.legend-actual,.legend-direct{background:var(--green)}.legend-target{background:var(--coral)}.legend-market{background:#74b596}.legend-retail{background:#c9dfd1}.table-toggle{margin-top:12px;border:0;background:transparent;color:#1f6b52;font-size:11px;font-weight:800;text-decoration:underline;cursor:pointer}.data-table{margin-top:10px;overflow:auto;border:1px solid var(--line);border-radius:12px}.data-table table{width:100%;border-collapse:collapse;font-size:11px}.data-table th,.data-table td{padding:8px 10px;border-bottom:1px solid #e5e8e5;text-align:left;white-space:nowrap}.data-table th{background:#edf3ef}.evidence-panel{position:fixed;right:24px;bottom:24px;z-index:20;width:min(390px,calc(100vw - 48px));border:1px solid #ffffff35;border-radius:18px;background:#13231eef;color:white;padding:20px;box-shadow:0 24px 70px #13231e45;backdrop-filter:blur(18px);opacity:0;visibility:hidden;transform:translateY(22px);pointer-events:none;transition:transform .22s,opacity .18s,visibility .18s}.evidence-panel[data-open="true"]{opacity:1;visibility:visible;transform:none;pointer-events:auto}.evidence-panel span{color:var(--lime);font-size:10px;font-weight:900;letter-spacing:.12em}.evidence-panel strong{display:block;margin-top:7px}.evidence-panel p{margin:7px 0 0;color:#bdcbc5;font:11px/1.6 ui-monospace,monospace;overflow-wrap:anywhere}.evidence-panel button{position:absolute;right:12px;top:10px;border:0;background:transparent;color:white;font-size:18px;cursor:pointer}.datum:focus,.line-dot:focus{outline:none;filter:drop-shadow(0 0 5px #c8ef72)}.candidate-note{max-width:1240px;margin:18px auto 0;border:1px dashed #13231e33;border-radius:16px;padding:14px 18px;color:#53665e;font-size:12px}.candidate-note strong{color:var(--ink)}body.ppt-capture{width:1280px;height:720px;overflow:hidden;background:var(--paper)}.ppt-capture .hero{width:1280px;height:720px;max-width:none;margin:0;border-radius:0;padding:72px 78px;display:flex;flex-direction:column;justify-content:center}.ppt-capture .chart-card{width:1280px;height:720px;border:0;border-radius:0;padding:42px 58px;display:flex;flex-direction:column}.ppt-capture .chart-heading h3{font-size:34px}.ppt-capture .chart-heading p{font-size:15px}.ppt-capture .chart-svg{flex:1;min-height:0;max-height:430px;margin:10px auto 4px}.ppt-capture .annotation{flex:none;padding:12px 16px}.ppt-capture .table-toggle,.ppt-capture .data-table{display:none}.ppt-capture .evidence-link{grid-column:2/4;grid-row:3;pointer-events:none;text-align:right;text-decoration:none;font-size:0}.ppt-capture .evidence-link:after{content:"Evidence · " attr(data-evidence);font-size:9px;color:#53665e}.ppt-capture .legend{margin-top:8px}@media(max-width:900px){.topbar nav{display:none}.hero{margin:12px;border-radius:22px;padding:42px 28px}.report{padding:0 12px 70px}.scenario{padding:22px 16px}.scenario-head{grid-template-columns:1fr}.metrics,.chart-grid{grid-template-columns:1fr}.chart-wide{grid-column:auto}.chart-heading{display:block}.unit-pill{display:inline-block;margin-top:10px}.annotation{grid-template-columns:1fr}.annotation>span,.annotation p,.evidence-link{grid-column:auto;grid-row:auto}}@media print{@page{size:A4 portrait;margin:11mm}.topbar,.table-toggle,.legend,.evidence-panel,.candidate-note{display:none!important}body{background:white}.hero{margin:0;border-radius:0;padding:34px 38px;min-height:275mm;display:flex;flex-direction:column;justify-content:center;break-after:page}.report{max-width:none;padding:0}.scenario{border:0;background:white;padding:0;break-before:page}.scenario-head{gap:10px}.metrics{gap:7px;margin:10px 0}.metric{padding:10px 12px}.metric strong{font-size:21px}.metric p{margin-top:3px}.chart-grid{display:block}.chart-card{margin-top:8px;break-inside:avoid;padding:13px 15px}.chart-card:nth-child(2){break-before:page;margin-top:55mm}.chart-svg{max-height:220px;margin:8px 0 2px}.annotation{padding:8px 10px}.annotation p{font-size:9px}.evidence-link{pointer-events:none;text-decoration:none;font-size:0}.evidence-link:after{content:"Evidence · " attr(data-evidence);font-size:7px;color:#53665e}.data-table{display:none}.scenario h2{font-size:27px}.chart-heading h3{font-size:17px}.chart-heading p{font-size:10px}.chart-kicker,.scenario-kicker{font-size:8px}.unit-pill{font-size:8px;padding:4px 7px}}
</style></head><body><header class="topbar"><div class="brand"><b>FORGE</b> · EVIDENCE STORY R0</div><nav><a href="#category">品类组合</a><a href="#trend">增长趋势</a><button class="ghost-btn" type="button" onclick="window.print()">打印候选</button></nav></header><section class="hero"><span class="eyebrow">CHART STORYTELLING · VISUAL GATE</span><h1>从图表堆砌，<br>到决策叙事。</h1><div class="rule"></div><p>每张图回答一个不同问题：谁领先、集中在哪里、何时失速、增长从哪里来。标注必须回到可复算 Evidence。</p><div class="hero-meta"><span>4 张互补视图</span><span>5 个证据绑定标注</span><span>HTML · PDF · PPTX 同源</span></div></section><aside class="candidate-note"><strong>R0 候选：</strong>此页面只验证 ChartArtifact v2 的信息架构、视觉与交互方向，不替换生产 Renderer，也不修改生产 Skills/Prompt。</aside><main class="report"><section id="category" class="scenario"><div class="scenario-head"><div><span class="scenario-kicker">SCENARIO A · PORTFOLIO</span><h2>品类组合：规模与集中度必须分开看</h2></div><span class="quality">QUALITY READY · 唯一品类粒度</span></div><div class="metrics"><article class="metric"><span>总销售额</span><strong>__CATEGORY_TOTAL__</strong><p>10 个唯一品类，可完整复算</p></article><article class="metric"><span>头部三类贡献</span><strong>45.8%</strong><p>不存在单一绝对赢家</p></article><article class="metric"><span>80% 覆盖所需</span><strong>6 类</strong><p>资源配置需要分层</p></article></div><div class="chart-grid">__CATEGORY_RANKING____CATEGORY_PARETO__</div></section><section id="trend" class="scenario"><div class="scenario-head"><div><span class="scenario-kicker">SCENARIO B · MOMENTUM</span><h2>月度增长：先识别拐点，再拆解贡献</h2></div><span class="quality">QUALITY READY · 连续月粒度</span></div><div class="metrics"><article class="metric"><span>六月销售额</span><strong>__JUNE_TOTAL__</strong><p>阶段新高</p></article><article class="metric"><span>六月超目标</span><strong>+9.1%</strong><p>五月起恢复到目标线上方</p></article><article class="metric"><span>四月至六月增量</span><strong>¥174K</strong><p>直营贡献其中 50%</p></article></div><div class="chart-grid">__LINE_CHART____AREA_CHART__</div></section></main><aside class="evidence-panel" id="evidence-panel" data-open="false" aria-live="polite"><button type="button" aria-label="关闭证据">×</button><span>EVIDENCE TRACE</span><strong id="evidence-note">选择标注查看证据</strong><p id="evidence-refs"></p></aside><script>
const panel=document.getElementById('evidence-panel');document.querySelectorAll('.evidence-link').forEach(button=>button.addEventListener('click',()=>{document.getElementById('evidence-note').textContent=button.dataset.note;document.getElementById('evidence-refs').textContent=button.dataset.evidence;panel.dataset.open='true'}));panel.querySelector('button').addEventListener('click',()=>panel.dataset.open='false');document.querySelectorAll('.table-toggle').forEach(button=>button.addEventListener('click',()=>{const table=button.nextElementSibling;const open=table.hasAttribute('hidden');table.toggleAttribute('hidden',!open);button.setAttribute('aria-expanded',String(open));button.textContent=open?'收起数据表':'查看数据表'}));document.querySelectorAll('[data-series]').forEach(button=>button.addEventListener('click',()=>{const on=button.getAttribute('aria-pressed')==='true';button.setAttribute('aria-pressed',String(!on));document.querySelectorAll('.series-'+button.dataset.series).forEach(node=>node.classList.toggle('is-muted',on))}));document.querySelectorAll('[data-area]').forEach(button=>button.addEventListener('click',()=>{document.querySelectorAll('[data-area]').forEach(item=>item.setAttribute('aria-pressed',String(item===button)));document.querySelectorAll('[data-area-series]').forEach(node=>node.classList.toggle('is-muted',node.dataset.areaSeries!==button.dataset.area))}));
</script></body></html>"""
    return (template.replace("__CATEGORY_TOTAL__",_compact(category_total)).replace("__JUNE_TOTAL__",_compact(june[4]))
            .replace("__CATEGORY_RANKING__",_category_ranking(category)).replace("__CATEGORY_PARETO__",_category_pareto(category))
            .replace("__LINE_CHART__",_line_chart(monthly)).replace("__AREA_CHART__",_stacked_area(monthly)))


def _pptx(category: dict[str, Any], monthly: dict[str, Any], target: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    prs=Presentation();prs.slide_width=Inches(13.333);prs.slide_height=Inches(7.5);blank=prs.slide_layouts[6]
    C={"ink":RGBColor(19,35,30),"paper":RGBColor(238,234,221),"card":RGBColor(255,253,247),"moss":RGBColor(18,63,49),"green":RGBColor(47,116,93),"mint":RGBColor(138,211,173),"lime":RGBColor(200,239,114),"coral":RGBColor(226,122,95),"line":RGBColor(215,221,215),"slate":RGBColor(113,128,121)}
    def bg(slide,color): slide.background.fill.solid();slide.background.fill.fore_color.rgb=color
    def box(slide,text,x,y,w,h,size=16,color=None,bold=False,font="Avenir Next",align=PP_ALIGN.LEFT):
        shape=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h));f=shape.text_frame;f.clear();f.word_wrap=True;p=f.paragraphs[0];p.alignment=align;r=p.add_run();r.text=str(text);r.font.name=font;r.font.size=Pt(size);r.font.bold=bold;r.font.color.rgb=color or C["ink"];return shape
    def header(slide,kicker,title,page):
        box(slide,kicker,0.75,0.42,5,0.25,9,C["green"],True);box(slide,title,0.75,0.76,10.8,0.6,27,C["ink"],True,"Georgia");box(slide,f"{page:02d}",12,0.48,.55,.22,9,C["slate"],True,align=PP_ALIGN.RIGHT)
        line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.75),Inches(1.42),Inches(11.8),Inches(.02));line.fill.solid();line.fill.fore_color.rgb=C["line"];line.line.fill.background()
    slide=prs.slides.add_slide(blank);bg(slide,C["moss"]);bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(.75),Inches(.75),Inches(.12),Inches(5.9));bar.fill.solid();bar.fill.fore_color.rgb=C["lime"];bar.line.fill.background();box(slide,"FORGE · EVIDENCE STORY R0",1.2,.88,5,.3,10,C["lime"],True);box(slide,"从图表堆砌，\n到决策叙事。",1.2,1.55,10.7,1.8,37,C["card"],True,"Georgia");box(slide,"每张图回答一个不同问题；每个标注都回到可复算 Evidence。",1.2,4.45,9.8,.8,20,RGBColor(220,235,228),False,"Georgia");box(slide,"ChartArtifact v2 · Visual Gate · 不进入生产主链",1.2,6.55,7,.25,9,C["lime"],True)
    rows=category["query_result"]["rows"];slide=prs.slides.add_slide(blank);bg(slide,C["paper"]);header(slide,"01 · CATEGORY RANKING","头部品类接近，资源不应单点押注",2);maxv=max(r[1] for r in rows);points=rows[:7]+[["其他 3 类",sum(r[1] for r in rows[7:]),0]]
    for i,(label,value,_o) in enumerate(points):
        y=1.72+i*.58;box(slide,label,.8,y,1.15,.25,10,C["slate"],True);w=6.35*value/maxv;shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(2.05),Inches(y),Inches(w),Inches(.28));shape.fill.solid();shape.fill.fore_color.rgb=C["lime"] if i==0 else C["coral"] if i==7 else C["green"];shape.line.fill.background();box(slide,_compact(value),2.15+w,y,.85,.25,9,C["ink"],True)
    box(slide,"“其他 3 类”为尾部合计",.8,6.62,3,.22,8,C["slate"])
    note=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(9.55),Inches(2.0),Inches(2.85),Inches(2.2));note.fill.solid();note.fill.fore_color.rgb=C["card"];note.line.color.rgb=C["line"];box(slide,"关键标注",9.87,2.25,1.8,.2,9,C["green"],True);box(slide,"第一名仅领先\n第二名 4.4%",9.87,2.66,2.15,.8,17,C["ink"],True,"Georgia");box(slide,"组合经营优于单点押注",9.87,3.58,2.05,.35,10,C["slate"])
    vals=[r[1] for r in rows];total=sum(vals);cum=[];run=0
    for v in vals:run+=v;cum.append(run/total)
    slide=prs.slides.add_slide(blank);bg(slide,C["paper"]);header(slide,"02 · CONTRIBUTION","经营六个品类，覆盖 82.2% 销售额",3);left,top,w,h=.95,1.9,10.9,3.9
    line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left),Inches(top+h*.2),Inches(w),Inches(.02));line.fill.solid();line.fill.fore_color.rgb=C["coral"];line.line.fill.background();box(slide,"80% COVERAGE",10.4,top+h*.2-.3,1.4,.2,8,C["coral"],True)
    prev=None
    for i,(row,share) in enumerate(zip(rows,cum)):
        x=left+i*w/(len(rows)-1);y=top+h*(1-share);bar_h=1.55*row[1]/max(vals);bar=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x-.22),Inches(6.0-bar_h),Inches(.44),Inches(bar_h));bar.fill.solid();bar.fill.fore_color.rgb=RGBColor(205,224,213);bar.line.fill.background();dot=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x-.06),Inches(y-.06),Inches(.12 if i!=5 else .2),Inches(.12 if i!=5 else .2));dot.fill.solid();dot.fill.fore_color.rgb=C["lime"] if i==5 else C["green"];dot.line.fill.background();
        if prev:
            px,py=prev;seg=slide.shapes.add_connector(1,Inches(px),Inches(py),Inches(x),Inches(y));seg.line.color.rgb=C["green"];seg.line.width=Pt(2)
        prev=(x,y);box(slide,row[0][:2],x-.18,6.15,.4,.2,8,C["slate"],align=PP_ALIGN.CENTER)
    for share in (.2,.4,.6,.8,1.0): box(slide,f"{share:.0%}",12.0,top+h*(1-share)-.1,.45,.18,8,C["slate"],align=PP_ALIGN.RIGHT)
    box(slide,"销售额（柱） · 累计占比（线）",.95,6.58,3.4,.2,8,C["slate"]);box(slide,"第六类跨过 80% 阈值",8.65,2.0,3,.4,15,C["ink"],True,"Georgia")
    mrows=monthly["query_result"]["rows"];slide=prs.slides.add_slide(blank);bg(slide,C["paper"]);header(slide,"03 · TREND","四月失速，五月完成反转",4);low,high=650000,930000;coords=[];targets=[]
    for i,r in enumerate(mrows):
        x=1+i*10.4/(len(mrows)-1);coords.append((x,5.8-3.7*(r[4]-low)/(high-low)));targets.append((x,5.8-3.7*(r[5]-low)/(high-low)));box(slide,r[0][5:]+"月",x-.25,6.12,.5,.2,9,C["slate"],align=PP_ALIGN.CENTER)
    for points,color,dash in [(targets,C["coral"],True),(coords,C["green"],False)]:
        for a,b in zip(points,points[1:]):seg=slide.shapes.add_connector(1,Inches(a[0]),Inches(a[1]),Inches(b[0]),Inches(b[1]));seg.line.color.rgb=color;seg.line.width=Pt(2.5);seg.line.dash_style=2 if dash else None
    for i,(x,y) in enumerate(coords):dot=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x-.07),Inches(y-.07),Inches(.14),Inches(.14));dot.fill.solid();dot.fill.fore_color.rgb=C["lime"] if i in (3,5) else C["green"];dot.line.fill.background()
    box(slide,"四月低于目标 6.4%",7.9,2.2,3,.45,17,C["coral"],True,"Georgia");box(slide,"六月较目标高 9.1%",9.15,4.25,2.6,.4,15,C["green"],True,"Georgia");box(slide,"Y 轴从 650K 起，用于观察目标偏差",.75,6.62,3.8,.2,8,C["slate"])
    slide=prs.slides.add_slide(blank);bg(slide,C["paper"]);header(slide,"04 · COMPOSITION","直营贡献了反转后 50% 的新增量",5);series=[("直营",1,C["green"]),("平台",2,C["mint"]),("门店",3,RGBColor(201,223,209))]
    for legend_index,(label,_idx,color) in enumerate(series):
        lx=.9+legend_index*1.05;dot=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(lx),Inches(1.56),Inches(.13),Inches(.13));dot.fill.solid();dot.fill.fore_color.rgb=color;dot.line.fill.background();box(slide,label,lx+.18,1.51,.65,.2,9,C["slate"],True)
    for i,r in enumerate(mrows):
        x=.85+i*1.65;bottom=6.0
        for label,idx,color in reversed(series):
            bh=3.7*r[idx]/930000;shape=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(bottom-bh),Inches(.9),Inches(bh));shape.fill.solid();shape.fill.fore_color.rgb=color;shape.line.fill.background();bottom-=bh
        box(slide,r[0][5:]+"月",x,6.16,.9,.2,9,C["slate"],align=PP_ALIGN.CENTER)
    box(slide,"APR → JUN",10.35,2.0,1.65,.25,9,C["green"],True);box(slide,"+174K",10.35,2.4,1.85,.55,26,C["ink"],True,"Georgia");box(slide,"直营贡献 87K\n占新增量 50%",10.35,3.15,1.85,.8,13,C["green"],True)
    prs.save(target)


def _pptx_from_images(images: list[Path], target: Path) -> None:
    """Build the final R0 deck from deterministic, full-slide candidate renders."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for image in images:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
    presentation.save(target)


def generate(output_dir: Path) -> dict[str, str]:
    output_dir.mkdir(parents=True,exist_ok=True)
    category=_load("category-comparison.json");monthly=_load("time-series.json")
    html_path=output_dir/"h5-chart-storytelling.html";html_path.write_text(_html(category,monthly))
    pptx_path=output_dir/"h5-chart-storytelling.pptx";_pptx(category,monthly,pptx_path)
    pdf_path=output_dir/"h5-chart-storytelling.pdf";png_path=output_dir/"h5-chart-storytelling.png"
    try:
        from playwright.sync_api import sync_playwright
        with sync_playwright() as p:
            browser=p.chromium.launch(headless=True);page=browser.new_page(viewport={"width":1600,"height":1000});errors=[]
            page.on("console",lambda message: errors.append(message.text) if message.type=="error" else None);page.on("pageerror",lambda error:errors.append(str(error)))
            page.goto(html_path.resolve().as_uri());page.wait_for_load_state("networkidle");page.screenshot(path=str(png_path),full_page=True);page.pdf(path=str(pdf_path),format="A4",print_background=True,margin={"top":"8mm","right":"8mm","bottom":"8mm","left":"8mm"})
            slide_dir=output_dir/"pptx-slides";slide_dir.mkdir(exist_ok=True);slide_paths=[]
            selectors=[".hero",*[f'.chart-card:nth-of-type({index})' for index in range(1,5)]]
            for index,selector in enumerate(selectors,1):
                page.goto(html_path.resolve().as_uri());page.wait_for_load_state("networkidle")
                if index == 1:
                    page.evaluate("""selector => { const source=document.querySelector(selector); document.body.innerHTML=source.outerHTML; document.body.className='ppt-capture'; }""",selector)
                else:
                    page.evaluate("""chartIndex => { const source=document.querySelectorAll('.chart-card')[chartIndex-1]; document.body.innerHTML=source.outerHTML; document.body.className='ppt-capture'; }""",index-1)
                page.set_viewport_size({"width":1280,"height":720});page.wait_for_timeout(50)
                slide_path=slide_dir/f"slide-{index:02d}.png";page.screenshot(path=str(slide_path),clip={"x":0,"y":0,"width":1280,"height":720});slide_paths.append(slide_path)
            _pptx_from_images(slide_paths,pptx_path);browser.close()
            if errors: raise RuntimeError(f"candidate browser errors: {errors}")
    except ImportError as error:
        raise RuntimeError("Playwright is required to generate the R0 HTML/PDF/PPTX evidence pack") from error
    outputs = {"html": html_path, "pdf": pdf_path, "pptx": pptx_path, "png": png_path}
    missing = [name for name, path in outputs.items() if not path.is_file()]
    if missing:
        raise RuntimeError(f"candidate generation did not produce: {', '.join(missing)}")
    return {name: str(path) for name, path in outputs.items()}


def main() -> None:
    parser=argparse.ArgumentParser();parser.add_argument("--output-dir",type=Path,required=True);args=parser.parse_args();print(json.dumps(generate(args.output_dir),ensure_ascii=False))


if __name__=="__main__":main()
