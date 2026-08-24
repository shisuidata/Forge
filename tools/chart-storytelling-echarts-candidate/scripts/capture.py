#!/usr/bin/env python3
"""Capture deterministic HTML/PDF/PPTX evidence from the focused candidate."""
from __future__ import annotations

import argparse
import json
import subprocess
import tempfile
from pathlib import Path
from urllib.parse import urlencode

from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches

SLIDES = ["cover", "ranking", "pareto", "trend", "contribution"]


def build_deck(images: list[Path], target: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for image in images:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
    presentation.save(target)


def capture(base_url: str, output_dir: Path) -> dict[str, object]:
    output_dir.mkdir(parents=True, exist_ok=True)
    slide_dir = output_dir / "slides"
    slide_dir.mkdir(exist_ok=True)
    pdf_path = output_dir / "echarts-focused-report.pdf"
    pptx_path = output_dir / "echarts-focused-report.pptx"
    errors: list[str] = []

    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1600, "height": 1000}, device_scale_factor=1)
        page.on("console", lambda message: errors.append(message.text) if message.type == "error" else None)
        page.on("pageerror", lambda error: errors.append(str(error)))
        page.goto(base_url, wait_until="networkidle")
        page.wait_for_selector('body[data-ready="true"]')
        page.pdf(
            path=str(pdf_path),
            format="A4",
            landscape=True,
            print_background=True,
            display_header_footer=False,
            margin={"top": "7mm", "right": "7mm", "bottom": "7mm", "left": "7mm"},
        )

        slide_paths: list[Path] = []
        for index, view in enumerate(SLIDES, start=1):
            slide_page = browser.new_page(viewport={"width": 1280, "height": 720}, device_scale_factor=1)
            slide_page.on("console", lambda message: errors.append(message.text) if message.type == "error" else None)
            slide_page.on("pageerror", lambda error: errors.append(str(error)))
            separator = "&" if "?" in base_url else "?"
            slide_page.goto(f"{base_url}{separator}{urlencode({'media': 'slide', 'view': view})}", wait_until="networkidle")
            slide_page.wait_for_selector('body[data-ready="true"]')
            slide_path = slide_dir / f"slide-{index:02d}-{view}.png"
            slide_page.screenshot(path=str(slide_path))
            with Image.open(slide_path) as image:
                assert image.size == (1280, 720)
            slide_paths.append(slide_path)
            slide_page.close()
        browser.close()

    assert errors == [], errors
    build_deck(slide_paths, pptx_path)

    extractor = Path(__file__).with_name("extract_pdf.swift")
    with tempfile.TemporaryDirectory(prefix="forge-pdf-gate-") as temporary_dir:
        binary = Path(temporary_dir) / "extract-pdf"
        subprocess.run(
            ["xcrun", "swiftc", "-framework", "PDFKit", str(extractor), "-o", str(binary)],
            check=True,
            capture_output=True,
            text=True,
        )
        extraction = subprocess.run(
            [str(binary), str(pdf_path)],
            check=True,
            capture_output=True,
            text=True,
        ).stdout
    first_line, _, pdf_text = extraction.partition("\n")
    assert first_line.startswith("PAGES="), first_line
    pdf_pages = int(first_line.removeprefix("PAGES="))
    assert pdf_pages == 5, f"expected 5 report pages, got {pdf_pages}"
    assert all(phrase in pdf_text for phrase in (
        "品类组合与增长诊断",
        "头部接近",
        "经营六个品类",
        "四月失速",
        "4→6 月新增 174K",
        "87K + 53K + 34K = 174K",
    ))
    assert not any(token in pdf_text for token in ("file://", "/Users/", "localhost", "127.0.0.1"))

    deck = Presentation(pptx_path)
    assert len(deck.slides) == 5
    assert all(len(slide.shapes) == 1 for slide in deck.slides)

    result = {
        "pdf": str(pdf_path),
        "pdf_bytes": pdf_path.stat().st_size,
        "pdf_pages": pdf_pages,
        "pptx": str(pptx_path),
        "pptx_bytes": pptx_path.stat().st_size,
        "pptx_slides": len(deck.slides),
        "slide_images": [str(path) for path in slide_paths],
        "errors": errors,
    }
    (output_dir / "static-gate.json").write_text(json.dumps(result, ensure_ascii=False, indent=2))
    return result


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--url", default="http://127.0.0.1:4175/")
    parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args()
    print(json.dumps(capture(args.url, args.output_dir), ensure_ascii=False))


if __name__ == "__main__":
    main()
