"""Immutable report publication and deterministic HTML/PDF/PPTX projection.

This service consumes already-approved Pi artifacts.  It never queries a data
source and never invokes an LLM.  SQLite stores metadata; large files live in a
mode-700 artifact directory.
"""
from __future__ import annotations

import hashlib
import html
import json
import os
import secrets
import shutil
import sqlite3
import subprocess
import tempfile
import threading
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from uuid import uuid4

_FORBIDDEN = ("<think", "system prompt", "tool call", "chain-of-thought", "api_key", "password", "secret")


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _canonical(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))


def _sha256(value: Any) -> str:
    return "sha256:" + hashlib.sha256(_canonical(value).encode()).hexdigest()


def _safe_text(value: Any, limit: int = 20_000) -> str:
    text = str(value or "").strip()[:limit]
    lowered = text.lower()
    if any(marker in lowered for marker in _FORBIDDEN):
        raise ValueError("report input contains forbidden reasoning, prompt, or secret material")
    return text


def _atomic_write(path: Path, data: bytes) -> None:
    path.parent.mkdir(parents=True, exist_ok=True, mode=0o700)
    fd, tmp_name = tempfile.mkstemp(prefix=f".{path.name}.", dir=path.parent)
    try:
        os.fchmod(fd, 0o600)
        with os.fdopen(fd, "wb") as stream:
            stream.write(data)
            stream.flush()
            os.fsync(stream.fileno())
        os.replace(tmp_name, path)
    finally:
        try:
            os.unlink(tmp_name)
        except FileNotFoundError:
            pass


def _svg_chart(chart: dict[str, Any], query: dict[str, Any]) -> str:
    columns = query.get("columns") or []
    rows = (query.get("rows") or [])[:12]
    dimension = chart.get("dimension")
    measures = chart.get("measures") or []
    if dimension not in columns or not measures or measures[0] not in columns:
        return ""
    di, mi = columns.index(dimension), columns.index(measures[0])
    points: list[tuple[str, float]] = []
    for row in rows:
        if not isinstance(row, list) or max(di, mi) >= len(row):
            continue
        try:
            value = float(row[mi])
        except (TypeError, ValueError):
            continue
        points.append((str(row[di]), value))
    if not points:
        return ""
    width, height, pad = 820, 320, 46
    max_value = max(abs(value) for _, value in points) or 1
    if chart.get("chart_type") == "line":
        step = (width - pad * 2) / max(1, len(points) - 1)
        coords = []
        labels = []
        for index, (label, value) in enumerate(points):
            x = pad + index * step
            y = height - pad - (max(0, value) / max_value) * (height - pad * 2)
            coords.append(f"{x:.1f},{y:.1f}")
            labels.append(f'<text x="{x:.1f}" y="{height-18}" text-anchor="middle">{html.escape(label[:12])}</text>')
        body = f'<polyline fill="none" stroke="#3d6b5d" stroke-width="4" points="{" ".join(coords)}"/>' + "".join(labels)
    else:
        gap = 8
        bar_width = max(12, (width - pad * 2) / len(points) - gap)
        chunks = []
        for index, (label, value) in enumerate(points):
            bar_height = (max(0, value) / max_value) * (height - pad * 2)
            x = pad + index * (bar_width + gap)
            y = height - pad - bar_height
            chunks.append(f'<rect x="{x:.1f}" y="{y:.1f}" width="{bar_width:.1f}" height="{bar_height:.1f}" rx="4" fill="#3d6b5d"/>')
            chunks.append(f'<text x="{x+bar_width/2:.1f}" y="{height-18}" text-anchor="middle">{html.escape(label[:12])}</text>')
        body = "".join(chunks)
    return f'<svg viewBox="0 0 {width} {height}" role="img" aria-label="{html.escape(str(chart.get("alt_text") or "图表"))}"><line x1="{pad}" y1="{height-pad}" x2="{width-pad}" y2="{height-pad}" stroke="#c8d0cc"/>{body}</svg>'


def _business_html(source: dict[str, Any]) -> str:
    report = source["business_report"]
    analysis = source["analysis"]
    query = source["query_result"]
    charts = source.get("charts") or []
    method = analysis.get("method_summary") or {}
    finding_source = report.get("key_findings") or analysis.get("findings") or []
    confidence_labels = {"high": "高置信", "medium": "中等置信", "low": "低置信"}
    priority_labels = {"high": "高优先级", "medium": "中优先级", "low": "低优先级"}

    finding_cards = []
    for index, item in enumerate(finding_source, 1):
        confidence = str(item.get("confidence") or "medium")
        statement = html.escape(_safe_text(item.get("statement")))
        interpretation = _safe_text(item.get("interpretation"))
        evidence_count = len(item.get("evidence_refs") or [])
        interpretation_html = f'<p class="finding-note">{html.escape(interpretation)}</p>' if interpretation else ""
        evidence_html = f'<span class="evidence-count">证据 {evidence_count} 条</span>' if evidence_count else ""
        finding_cards.append(
            f'<article class="finding-card"><div class="finding-index">{index:02d}</div>'
            f'<div class="finding-content"><div class="finding-meta"><span class="badge confidence-{confidence}">'
            f'{confidence_labels.get(confidence, "置信度待确认")}</span>{evidence_html}</div>'
            f'<h3>{statement}</h3>{interpretation_html}</div></article>'
        )

    recommendation_cards = []
    for item in report.get("recommendations") or []:
        priority = str(item.get("priority") or "medium")
        recommendation_cards.append(
            f'<article class="action-card"><span class="badge priority-{priority}">'
            f'{priority_labels.get(priority, "优先级待确认")}</span>'
            f'<h3>{html.escape(_safe_text(item.get("action")))}</h3>'
            f'<p>{html.escape(_safe_text(item.get("rationale")))}</p></article>'
        )

    limitations_source = report.get("limitations") or analysis.get("limitations") or []
    limitations = "".join(f"<li>{html.escape(_safe_text(item))}</li>" for item in limitations_source)
    next_steps = "".join(f"<li>{html.escape(_safe_text(item))}</li>" for item in report.get("next_steps") or [])
    steps = "".join(
        f'<li><span>{index:02d}</span>{html.escape(_safe_text(item))}</li>'
        for index, item in enumerate(method.get("approach_steps") or [], 1)
    )
    chart_html = "".join(
        f'<section class="report-section chart"><div class="section-heading"><p>DATA VIEW</p>'
        f'<h2>{html.escape(_safe_text(chart.get("title")))}</h2></div>{_svg_chart(chart, query)}</section>'
        for chart in charts
    )
    table_columns = query.get("columns") or []
    table_rows = (query.get("rows") or [])[:20]
    head = "".join(f"<th>{html.escape(str(column))}</th>" for column in table_columns)
    body = "".join(
        "<tr>" + "".join(f"<td>{html.escape(str(value))}</td>" for value in row) + "</tr>"
        for row in table_rows
    )
    title = html.escape(_safe_text(report.get("title") or source.get("title") or "分析报告"))
    summary = html.escape(_safe_text(report.get("executive_summary") or analysis.get("summary")))
    objective = html.escape(_safe_text(method.get("objective")))
    dimensions = html.escape("、".join(method.get("dimensions") or []))
    baseline = html.escape(_safe_text(method.get("comparison_baseline")))
    report_id = html.escape(_safe_text(source.get("report_id"), 128))
    findings_html = "".join(finding_cards) or '<p class="empty-state">当前报告没有可发布的关键发现。</p>'
    recommendations_html = "".join(recommendation_cards) or '<p class="empty-state">当前没有已确认的行动建议。</p>'
    limitation_html = limitations or "<li>当前没有额外限制说明。</li>"
    next_steps_section = (
        f'<section class="report-section next-steps"><div class="section-heading"><p>NEXT</p><h2>下一步</h2></div><ol>{next_steps}</ol></section>'
        if next_steps else ""
    )
    return f"""<!doctype html><html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>{title}</title><style>
:root{{--ink:#17201d;--muted:#65716c;--paper:#f1eee5;--card:#fffefa;--moss:#194d3a;--accent:#80cfa8;--lime:#c9ef78;--amber:#c9831f;--coral:#bd5b43;--line:#d8ded9}}*{{box-sizing:border-box}}html{{scroll-behavior:smooth}}body{{margin:0;background:linear-gradient(135deg,rgba(25,77,58,.035) 25%,transparent 25%) 0 0/20px 20px,var(--paper);color:var(--ink);font:16px/1.72 ui-sans-serif,-apple-system,BlinkMacSystemFont,"Segoe UI",sans-serif}}main{{max-width:1120px;margin:auto;padding:64px 28px 104px}}.toolbar{{position:sticky;top:0;z-index:3;display:flex;gap:10px;flex-wrap:wrap;align-items:center;padding:12px 24px;background:#f1eee5ed;backdrop-filter:blur(14px);border-bottom:1px solid var(--line)}}.toolbar a,.toolbar button{{border:1px solid #aab8b1;background:var(--card);color:var(--ink);border-radius:999px;padding:8px 14px;text-decoration:none;cursor:pointer;font:inherit;font-size:13px;font-weight:600;line-height:1.2}}.toolbar a:hover,.toolbar button:hover{{border-color:var(--moss);color:var(--moss)}}#share-result{{flex-basis:100%;font-size:13px;color:var(--muted);overflow-wrap:anywhere}}.report-hero{{position:relative;overflow:hidden;border-radius:28px;background:var(--moss);padding:54px clamp(26px,6vw,72px);color:#f6fbf8;box-shadow:0 28px 70px #17201d24}}.report-hero:after{{content:"";position:absolute;right:-90px;top:-130px;width:360px;height:360px;border:1px solid #c9ef7855;border-radius:50%;box-shadow:0 0 0 52px #c9ef780a,0 0 0 104px #c9ef7808}}.eyebrow,.section-heading p{{margin:0 0 10px;color:#92cfb5;font-size:12px;font-weight:800;letter-spacing:.18em}}h1{{position:relative;z-index:1;max-width:880px;margin:0;font:700 clamp(38px,6vw,72px)/1.06 ui-serif,Georgia,serif;letter-spacing:-.035em}}.hero-rule{{position:relative;z-index:1;width:72px;height:5px;margin:28px 0;background:var(--lime)}}.executive-brief{{position:relative;z-index:1;max-width:820px;border-left:3px solid var(--lime);padding:2px 0 2px 20px}}.executive-brief span{{display:block;margin-bottom:6px;color:#c9ef78;font-size:11px;font-weight:800;letter-spacing:.14em}}.executive-brief p{{margin:0;color:#eef7f2;font:20px/1.65 ui-serif,Georgia,serif}}.report-section{{margin-top:26px;border:1px solid var(--line);border-radius:22px;background:var(--card);padding:30px;box-shadow:0 12px 38px #23352e0d}}.section-heading{{display:flex;align-items:end;justify-content:space-between;gap:20px;margin-bottom:24px;border-bottom:1px solid #e5e9e6;padding-bottom:15px}}.section-heading p{{margin:0;color:#56806f}}.section-heading h2{{margin:0;font:700 clamp(24px,3vw,34px)/1.1 ui-serif,Georgia,serif;letter-spacing:-.02em}}.method-grid{{display:grid;grid-template-columns:minmax(0,.9fr) minmax(0,1.1fr);gap:30px}}.method-facts{{display:grid;gap:12px;margin:0}}.method-facts div{{border-left:3px solid #bdd8cb;padding:5px 0 5px 14px}}.method-facts dt{{color:var(--muted);font-size:11px;font-weight:800;letter-spacing:.1em}}.method-facts dd{{margin:3px 0 0;font-weight:650}}.method-steps{{list-style:none;margin:0;padding:0}}.method-steps li{{display:grid;grid-template-columns:34px 1fr;gap:12px;align-items:start;padding:10px 0;border-bottom:1px solid #edf0ee}}.method-steps li:last-child{{border-bottom:0}}.method-steps span,.finding-index{{color:#4f826f;font:700 13px/1.5 ui-monospace,monospace}}.findings-list{{display:grid;gap:14px}}.finding-card{{display:grid;grid-template-columns:48px 1fr;gap:16px;border:1px solid #dce5e0;border-radius:16px;padding:20px;background:#fbfdfb}}.finding-index{{padding-top:4px;font-size:18px}}.finding-meta{{display:flex;align-items:center;gap:9px;flex-wrap:wrap}}.badge{{display:inline-flex;border-radius:999px;padding:4px 9px;font-size:11px;font-weight:800;letter-spacing:.04em}}.confidence-high{{background:#e5f5eb;color:#17603f}}.confidence-medium{{background:#fff3d7;color:#875b11}}.confidence-low{{background:#ffebe6;color:#93442f}}.evidence-count{{color:var(--muted);font-size:12px}}.finding-card h3,.action-card h3{{margin:9px 0 0;font-size:17px;line-height:1.55}}.finding-note,.action-card p{{margin:8px 0 0;color:var(--muted)}}.action-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(250px,1fr));gap:14px}}.action-card{{border:1px solid #dce5e0;border-top:4px solid var(--accent);border-radius:14px;padding:19px;background:#fbfdfb}}.priority-high{{background:#fff0cf;color:#855208}}.priority-medium{{background:#eaf3ef;color:#315f50}}.priority-low{{background:#eff1f0;color:#66736e}}.risk-panel{{border-color:#e7c6bd;background:#fff6f2}}.risk-panel .section-heading{{border-color:#eed8d1}}.risk-panel .section-heading p{{color:var(--coral)}}.risk-list{{margin:0;padding-left:1.25rem}}.risk-list li::marker{{color:var(--coral)}}.risk-list li+li{{margin-top:10px}}.next-steps ol{{margin:0;padding-left:1.3rem}}.next-steps li+li{{margin-top:9px}}.chart svg{{display:block;width:100%;height:auto}}.chart svg text{{font-size:12px;fill:var(--muted)}}.table-wrap{{overflow:auto;border:1px solid #e1e6e3;border-radius:13px}}table{{width:100%;border-collapse:collapse;font-size:14px}}th,td{{text-align:left;border-bottom:1px solid #e5e9e6;padding:11px 13px;white-space:nowrap}}th{{position:sticky;top:0;background:#eef4f1;color:#315346;font-size:12px;letter-spacing:.03em}}tbody tr:last-child td{{border-bottom:0}}.empty-state{{color:var(--muted);font-style:italic}}footer{{margin-top:48px;border-top:1px solid #ccd6d1;padding-top:18px;color:var(--muted);font-size:13px}}@media(max-width:720px){{main{{padding:26px 14px 72px}}.toolbar{{padding:10px 12px}}.report-hero{{border-radius:20px;padding:38px 24px}}.executive-brief p{{font-size:17px}}.report-section{{padding:22px 18px}}.method-grid{{grid-template-columns:1fr}}.finding-card{{grid-template-columns:36px 1fr;padding:16px 14px}}.section-heading{{align-items:start;flex-direction:column;gap:4px}}}}@media(prefers-reduced-motion:reduce){{html{{scroll-behavior:auto}}}}@media print{{@page{{margin:14mm}}body{{background:#fff;font-size:11pt}}.toolbar{{display:none}}main{{max-width:none;padding:0}}.report-hero{{box-shadow:none;border-radius:0;padding:34px 38px;-webkit-print-color-adjust:exact;print-color-adjust:exact}}.report-section{{box-shadow:none;margin-top:14px;padding:20px}}.report-hero,.method,.chart,.risk-panel,.next-steps{{break-inside:avoid}}.section-heading{{break-after:avoid}}.finding-card,.action-card,.risk-panel,tr{{break-inside:avoid;-webkit-print-color-adjust:exact;print-color-adjust:exact}}thead{{display:table-header-group}}h1{{font-size:34pt}}.executive-brief p{{font-size:14pt}}}}
</style></head><body><nav class="toolbar" aria-label="报告操作"><a href="/reports/{report_id}/download/pdf">下载 PDF</a><a href="/reports/{report_id}/download/pptx">下载 PPTX</a><a href="/reports/{report_id}/technical">技术报告</a><button id="share">创建 7 天分享链接</button><span id="share-result" aria-live="polite"></span></nav><main><header class="report-hero"><p class="eyebrow">FORGE · EVIDENCE-BOUND ANALYSIS</p><h1>{title}</h1><div class="hero-rule"></div><div class="executive-brief"><span>EXECUTIVE SUMMARY</span><p>{summary}</p></div></header><section class="report-section method"><div class="section-heading"><p>METHOD</p><h2>分析思路</h2></div><div class="method-grid"><dl class="method-facts"><div><dt>目标</dt><dd>{objective}</dd></div><div><dt>观察维度</dt><dd>{dimensions}</dd></div><div><dt>对比基线</dt><dd>{baseline}</dd></div></dl><ol class="method-steps">{steps}</ol></div></section><section class="report-section findings"><div class="section-heading"><p>FINDINGS</p><h2>关键发现</h2></div><div class="findings-list">{findings_html}</div></section>{chart_html}<section class="report-section"><div class="section-heading"><p>EVIDENCE</p><h2>数据明细</h2></div><div class="table-wrap"><table><thead><tr>{head}</tr></thead><tbody>{body}</tbody></table></div></section><section class="report-section"><div class="section-heading"><p>ACTIONS</p><h2>建议行动</h2></div><div class="action-grid">{recommendations_html}</div></section>{next_steps_section}<section class="report-section risk-panel"><div class="section-heading"><p>LIMITATIONS</p><h2>限制与风险</h2></div><ul class="risk-list">{limitation_html}</ul></section><footer>本报告固定到不可变数据与分析版本；网页、PDF 与 PPTX 保持相同的信息优先级，不改变证据和结论边界。</footer></main><script>document.getElementById('share').onclick=async()=>{{const expires_at=new Date(Date.now()+7*86400000).toISOString();const r=await fetch('/api/reports/{report_id}/shares',{{method:'POST',headers:{{'content-type':'application/json'}},body:JSON.stringify({{expires_at}})}});const out=document.getElementById('share-result');if(!r.ok){{out.textContent='创建失败，请确认已登录且有权限';return;}}const x=await r.json();out.textContent=x.exchange_url;try{{await navigator.clipboard.writeText(x.exchange_url);out.textContent='分享链接已复制：'+x.exchange_url;}}catch{{}}}};</script></body></html>"""


def _technical_html(source: dict[str, Any]) -> str:
    tech = source["technical_report"]
    decision_rows = "".join(
        f"<tr><td>{html.escape(_safe_text(item.get('stage')))}</td><td>{html.escape(_safe_text(item.get('decision')))}</td><td>{html.escape(_safe_text(item.get('rationale')))}</td></tr>"
        for item in tech.get("decision_log") or []
    )
    lineage_rows = "".join(
        f"<tr><th>{html.escape(str(key))}</th><td>{html.escape(str(value))}</td></tr>"
        for key, value in (tech.get("lineage") or {}).items()
    )
    title = html.escape(_safe_text(tech.get("title")))
    return f"""<!doctype html><html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>{title}</title><style>:root{{--ink:#18201d;--muted:#66736e;--paper:#f5f6f5;--line:#ccd4d0;--code:#111a17}}*{{box-sizing:border-box}}body{{max-width:1040px;margin:0 auto;padding:52px 24px 90px;background:#fff;color:var(--ink);font:14px/1.68 ui-monospace,SFMono-Regular,Menlo,monospace;overflow-wrap:anywhere}}header{{border-bottom:2px solid #315f50;padding-bottom:22px;margin-bottom:36px}}h1{{margin:0;font:700 clamp(28px,5vw,46px)/1.12 ui-sans-serif,system-ui,sans-serif;letter-spacing:-.02em}}header p{{max-width:820px;color:var(--muted)}}h2{{margin:42px 0 12px;font:700 21px/1.2 ui-sans-serif,system-ui,sans-serif}}pre{{white-space:pre-wrap;word-break:break-word;background:var(--code);color:#d9eee5;padding:20px;border-radius:12px;overflow:auto}}table{{width:100%;border-collapse:collapse;table-layout:fixed}}th,td{{border:1px solid var(--line);padding:10px;text-align:left;vertical-align:top;word-break:break-word}}th{{width:24%;background:var(--paper);font-weight:700}}@media(max-width:640px){{body{{padding:28px 14px 60px}}table{{font-size:12px}}}}@media print{{body{{max-width:none;padding:0}}pre,table{{break-inside:avoid}}}}</style></head><body><header><h1>{title}</h1><p>该文档记录可复现的结构化决策、SQL、审批、执行和版本 lineage；不包含模型 hidden chain-of-thought、Prompt 或 Secret。</p></header><h2>SQL</h2><pre>{html.escape(_safe_text(tech.get('sql'), 100_000))}</pre><h2>审批与执行</h2><pre>{html.escape(_canonical({'approval': tech.get('approval'), 'execution': tech.get('execution')}))}</pre><h2>Decision Log</h2><table><tr><th>Stage</th><th>Decision</th><th>Rationale</th></tr>{decision_rows}</table><h2>Lineage</h2><table>{lineage_rows}</table></body></html>"""


class ReportStore:
    def __init__(self, db_path: str, artifact_dir: str):
        self.db_path = Path(db_path)
        self.artifact_dir = Path(artifact_dir)
        self.db_path.parent.mkdir(parents=True, exist_ok=True, mode=0o700)
        self.artifact_dir.mkdir(parents=True, exist_ok=True, mode=0o700)
        os.chmod(self.artifact_dir, 0o700)
        self._lock = threading.RLock()
        self._init()

    def _connect(self) -> sqlite3.Connection:
        conn = sqlite3.connect(self.db_path, timeout=5)
        conn.row_factory = sqlite3.Row
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("PRAGMA foreign_keys=ON")
        return conn

    def _init(self) -> None:
        with self._connect() as conn:
            conn.executescript("""
            CREATE TABLE IF NOT EXISTS reports (
              report_id TEXT PRIMARY KEY, task_run_id TEXT NOT NULL UNIQUE,
              org_id TEXT NOT NULL, team_id TEXT NOT NULL, user_id TEXT NOT NULL,
              revision INTEGER NOT NULL, bundle_hash TEXT NOT NULL, title TEXT NOT NULL,
              status TEXT NOT NULL, pdf_status TEXT NOT NULL, pptx_status TEXT NOT NULL,
              error TEXT, created_at TEXT NOT NULL, updated_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS report_attempts (
              attempt_id TEXT PRIMARY KEY, report_id TEXT NOT NULL, stage TEXT NOT NULL,
              status TEXT NOT NULL, started_at TEXT NOT NULL, finished_at TEXT, error TEXT
            );
            CREATE INDEX IF NOT EXISTS idx_report_attempts_report ON report_attempts(report_id, started_at);
            CREATE TABLE IF NOT EXISTS report_shares (
              share_id TEXT PRIMARY KEY, report_id TEXT NOT NULL, token_hash TEXT NOT NULL UNIQUE,
              scope TEXT NOT NULL CHECK(scope IN ('business')),
              expires_at TEXT NOT NULL, revoked_at TEXT, created_at TEXT NOT NULL
            );
            CREATE INDEX IF NOT EXISTS idx_report_shares_report ON report_shares(report_id, expires_at);
            CREATE TABLE IF NOT EXISTS report_download_audit (
              audit_id TEXT PRIMARY KEY, report_id TEXT NOT NULL, format TEXT NOT NULL,
              actor TEXT NOT NULL, created_at TEXT NOT NULL
            );
            """)
            now = _now()
            conn.execute(
                "UPDATE report_attempts SET status='interrupted', finished_at=?, error='process restarted; not replayed' WHERE status='running'",
                (now,),
            )
            conn.execute(
                "UPDATE reports SET status='failed', error='publication interrupted; explicit retry required', updated_at=? WHERE status='publishing'",
                (now,),
            )
        os.chmod(self.db_path, 0o600)

    def create(self, source: dict[str, Any]) -> dict[str, Any]:
        allowed = {
            "task_run_id", "org_id", "team_id", "user_id", "report_id", "revision",
            "bundle_hash", "title", "business_report", "analysis", "query_result", "charts",
            "technical_report",
        }
        if set(source) != allowed:
            raise ValueError("report input fields do not match the fixed contract")
        for key in ("task_run_id", "org_id", "team_id", "user_id", "report_id", "bundle_hash", "title"):
            if not isinstance(source.get(key), str) or not source[key]:
                raise ValueError(f"{key} is required")
        if not str(source["bundle_hash"]).startswith("sha256:") or len(source["bundle_hash"]) != 71:
            raise ValueError("bundle_hash is invalid")
        if not all(isinstance(source.get(key), dict) for key in ("business_report", "analysis", "query_result", "technical_report")):
            raise ValueError("report artifacts must be objects")
        if not isinstance(source.get("charts"), list) or len(source["charts"]) > 12:
            raise ValueError("charts must be a bounded list")
        query = source["query_result"]
        if not isinstance(query.get("columns"), list) or len(query["columns"]) > 200 or not isinstance(query.get("rows"), list) or len(query["rows"]) > 200:
            raise ValueError("query result projection exceeds report bounds")
        serialized = _canonical(source)
        if len(serialized.encode()) > 2_000_000:
            raise ValueError("report input exceeds 2 MB")
        _safe_text(serialized, 2_000_000)
        now = _now()
        with self._lock, self._connect() as conn:
            existing = conn.execute("SELECT * FROM reports WHERE task_run_id = ?", (source["task_run_id"],)).fetchone()
            if existing:
                if existing["bundle_hash"] != source["bundle_hash"]:
                    raise ValueError("task already published with a different bundle")
                return self._row(existing)
            conn.execute(
                "INSERT INTO reports VALUES (?, ?, ?, ?, ?, ?, ?, ?, 'publishing', 'pending', 'pending', NULL, ?, ?)",
                (source["report_id"], source["task_run_id"], source["org_id"], source["team_id"], source["user_id"],
                 int(source.get("revision") or 1), source["bundle_hash"], source["title"], now, now),
            )
            report_dir = self.artifact_dir / source["report_id"] / f"v{int(source.get('revision') or 1)}"
            _atomic_write(report_dir / "source.json", _canonical(source).encode())
        return self.get(source["report_id"])

    def _attempt(self, report_id: str, stage: str, operation):
        attempt_id, started = f"rpa_{uuid4().hex}", _now()
        with self._connect() as conn:
            conn.execute("INSERT INTO report_attempts VALUES (?, ?, ?, 'running', ?, NULL, NULL)",
                         (attempt_id, report_id, stage, started))
        try:
            result = operation()
            status, error = "succeeded", None
            return result
        except Exception as exc:
            status, error = "failed", str(exc)[:1000]
            raise
        finally:
            with self._connect() as conn:
                conn.execute("UPDATE report_attempts SET status=?, finished_at=?, error=? WHERE attempt_id=?",
                             (status, _now(), error, attempt_id))

    def build(self, report_id: str) -> dict[str, Any]:
        report = self.get(report_id)
        report_dir = self.artifact_dir / report_id / f"v{report['revision']}"
        source = json.loads((report_dir / "source.json").read_text())
        try:
            def html_job():
                _atomic_write(report_dir / "index.html", _business_html(source).encode())
                _atomic_write(report_dir / "technical.html", _technical_html(source).encode())
            self._attempt(report_id, "html", html_job)
            status, error = "published", None
        except Exception as exc:
            status, error = "failed", str(exc)[:1000]
        pdf_status = self._run_export(report_id, "pdf", report_dir, source) if status == "published" else "failed"
        pptx_status = self._run_export(report_id, "pptx", report_dir, source) if status == "published" else "failed"
        with self._lock, self._connect() as conn:
            conn.execute(
                "UPDATE reports SET status=?, pdf_status=?, pptx_status=?, error=?, updated_at=? WHERE report_id=?",
                (status, pdf_status, pptx_status, error, _now(), report_id),
            )
        return self.get(report_id)

    def _run_export(self, report_id: str, fmt: str, report_dir: Path, source: dict[str, Any]) -> str:
        try:
            def export_job():
                result = self._build_pdf(report_dir) if fmt == "pdf" else self._build_pptx(report_dir, source)
                if result != "ready":
                    raise RuntimeError(f"{fmt} exporter is unavailable or failed")
                return result
            self._attempt(report_id, fmt, export_job)
            return "ready"
        except Exception:
            return "failed"

    def retry_export(self, report_id: str, fmt: str) -> dict[str, Any]:
        if fmt not in {"pdf", "pptx"}:
            raise ValueError("unsupported export format")
        report = self.get(report_id)
        if report["status"] != "published":
            raise ValueError("report HTML is not published")
        report_dir = self.artifact_dir / report_id / f"v{report['revision']}"
        source = json.loads((report_dir / "source.json").read_text())
        result = self._run_export(report_id, fmt, report_dir, source)
        with self._connect() as conn:
            conn.execute(f"UPDATE reports SET {fmt}_status=?, updated_at=? WHERE report_id=?", (result, _now(), report_id))
        return self.get(report_id)

    def _build_pdf(self, report_dir: Path) -> str:
        chrome = next((path for path in (shutil.which("google-chrome"), shutil.which("chromium"), shutil.which("chromium-browser")) if path), None)
        if chrome is None:
            return "failed"
        target = (report_dir / "report.pdf").resolve()
        command = [chrome, "--headless", "--disable-gpu", "--no-sandbox", f"--print-to-pdf={target}", (report_dir / "index.html").resolve().as_uri()]
        result = subprocess.run(command, capture_output=True, timeout=90, check=False)
        if result.returncode != 0 or not target.exists():
            return "failed"
        os.chmod(target, 0o600)
        return "ready"

    def _build_pptx(self, report_dir: Path, source: dict[str, Any]) -> str:
        try:
            from textwrap import wrap

            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError:
            return "failed"

        prs = Presentation()
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
        blank = prs.slide_layouts[6]
        colors = {
            "ink": RGBColor(23, 32, 29), "muted": RGBColor(101, 113, 108),
            "paper": RGBColor(244, 241, 232), "card": RGBColor(255, 254, 250),
            "moss": RGBColor(25, 77, 58), "lime": RGBColor(201, 239, 120),
            "green": RGBColor(61, 107, 93), "amber": RGBColor(193, 122, 23),
            "coral": RGBColor(189, 91, 67), "line": RGBColor(216, 222, 217),
        }

        def background(slide, color):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def text_box(slide, text, left, top, width, height, *, size=18, color=None,
                     bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            frame = box.text_frame
            frame.clear(); frame.word_wrap = True; frame.vertical_anchor = valign
            frame.margin_left = frame.margin_right = Inches(0.02)
            frame.margin_top = frame.margin_bottom = Inches(0.02)
            paragraph = frame.paragraphs[0]
            paragraph.alignment = align
            run = paragraph.add_run(); run.text = _safe_text(text)
            run.font.name = font; run.font.size = Pt(size); run.font.bold = bold
            run.font.color.rgb = color or colors["ink"]
            return box

        def slide_header(slide, kicker, title, page_number, tone="green"):
            text_box(slide, kicker, 0.78, 0.48, 4.5, 0.3, size=10, bold=True,
                     color=colors[tone] if tone in colors else colors["green"])
            text_box(slide, title, 0.78, 0.83, 10.9, 0.62, size=28, bold=True, font="Georgia")
            text_box(slide, f"{page_number:02d}", 11.85, 0.58, 0.68, 0.32, size=10,
                     bold=True, color=colors["muted"], align=PP_ALIGN.RIGHT)
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(1.48), Inches(11.75), Inches(0.025))
            line.fill.solid(); line.fill.fore_color.rgb = colors["line"]; line.line.fill.background()

        def footer(slide):
            text_box(slide, "FORGE · EVIDENCE-BOUND REPORT", 0.78, 7.05, 5.5, 0.2,
                     size=8, bold=True, color=colors["muted"])

        def fragments(value, width=108):
            normalized = " ".join(_safe_text(value).split())
            return wrap(normalized, width=width, break_long_words=True, break_on_hyphens=False) or [""]

        def expand_records(records):
            expanded = []
            for label, value, badge in records:
                parts = fragments(value)
                for index, part in enumerate(parts):
                    expanded.append((label if index == 0 else "续", part, badge if index == 0 else ""))
            return expanded

        def card_slides(kicker, title, records, tone="green"):
            expanded = expand_records(records)
            if not expanded:
                expanded = [("—", "当前没有可发布内容。", "")]
            for offset in range(0, len(expanded), 3):
                page_records = expanded[offset:offset + 3]
                slide = prs.slides.add_slide(blank); background(slide, colors["paper"])
                slide_header(slide, kicker, title, len(prs.slides), tone)
                for index, (label, value, badge) in enumerate(page_records):
                    top = 1.78 + index * 1.68
                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(top), Inches(11.75), Inches(1.46)
                    )
                    card.fill.solid(); card.fill.fore_color.rgb = colors["card"]
                    card.line.color.rgb = colors["line"]
                    text_box(slide, label, 1.03, top + 0.22, 0.58, 0.32, size=12, bold=True,
                             color=colors[tone] if tone in colors else colors["green"])
                    if badge:
                        badge_color = (
                            colors["coral"] if "低置信" in badge or "需关注" in badge
                            else colors["amber"] if "中等置信" in badge or "高优先级" in badge
                            else colors["muted"] if "低优先级" in badge
                            else colors["green"]
                        )
                        badge_box = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.68), Inches(top + 0.15), Inches(1.55), Inches(0.32)
                        )
                        badge_box.fill.solid(); badge_box.fill.fore_color.rgb = colors["paper"]
                        badge_box.line.fill.background()
                        text_box(slide, badge, 1.78, top + 0.205, 1.35, 0.2, size=8, bold=True,
                                 color=badge_color)
                        value_top = top + 0.55
                    else:
                        value_top = top + 0.31
                    box = text_box(slide, value, 1.68, value_top, 10.35, 0.78, size=17)
                    box.text_frame.auto_size = None
                footer(slide)

        report = source["business_report"]
        analysis = source["analysis"]
        method = analysis.get("method_summary") or {}
        title = _safe_text(report.get("title") or source.get("title"))
        summary = _safe_text(report.get("executive_summary") or analysis.get("summary"))
        title_parts = fragments(title, 80)
        cover_title = title_parts[0] + ("…" if len(title_parts) > 1 else "")

        slide = prs.slides.add_slide(blank); background(slide, colors["moss"])
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(0.78), Inches(0.12), Inches(5.9))
        accent.fill.solid(); accent.fill.fore_color.rgb = colors["lime"]; accent.line.fill.background()
        text_box(slide, "FORGE ANALYSIS REPORT", 1.18, 0.88, 5.2, 0.35, size=11, bold=True, color=colors["lime"])
        text_box(slide, cover_title, 1.18, 1.45, 10.8, 2.15, size=38, bold=True, color=colors["card"], font="Georgia")
        summary_preview = fragments(summary, 150)[0]
        text_box(slide, summary_preview, 1.18, 4.28, 9.7, 1.35, size=19, color=RGBColor(224, 238, 231), font="Georgia")
        text_box(slide, "可信数据任务 · 不可变证据版本", 1.18, 6.55, 5.5, 0.25, size=9, bold=True, color=colors["lime"])

        if len(title_parts) > 1:
            card_slides("REPORT TITLE", "完整标题", [("标题", title, "")])
        card_slides("EXECUTIVE SUMMARY", "执行摘要", [("摘要", summary, "核心说明")])
        method_records = [
            ("目标", _safe_text(method.get("objective")), "METHOD"),
            ("维度", "、".join(method.get("dimensions") or []), "SCOPE"),
            ("基线", _safe_text(method.get("comparison_baseline")), "BASELINE"),
            *[(f"步骤 {index:02d}", _safe_text(item), "APPROACH")
              for index, item in enumerate(method.get("approach_steps") or [], 1)],
        ]
        card_slides("METHOD", "分析思路", method_records)

        confidence_labels = {"high": "高置信", "medium": "中等置信", "low": "低置信"}
        finding_source = report.get("key_findings") or analysis.get("findings") or []
        finding_records = []
        for index, finding in enumerate(finding_source, 1):
            statement = _safe_text(finding.get("statement"))
            interpretation = _safe_text(finding.get("interpretation"))
            combined = f"{statement} — {interpretation}" if interpretation else statement
            finding_records.append((f"{index:02d}", combined, confidence_labels.get(finding.get("confidence"), "置信度待确认")))
        card_slides("FINDINGS", "关键发现", finding_records)

        query = source.get("query_result") or {}
        columns, rows = query.get("columns") or [], (query.get("rows") or [])[:10]
        for chart in source.get("charts") or []:
            dimension, measures = chart.get("dimension"), chart.get("measures") or []
            if dimension not in columns or not measures or measures[0] not in columns:
                continue
            di, mi = columns.index(dimension), columns.index(measures[0])
            points = []
            for row in rows:
                try:
                    points.append((str(row[di]), max(0.0, float(row[mi]))))
                except (TypeError, ValueError, IndexError):
                    pass
            if not points:
                continue
            chart_title = _safe_text(chart.get("title"))
            chart_title_parts = fragments(chart_title, 70)
            if len(chart_title_parts) > 1:
                card_slides("DATA VIEW", "完整图表标题", [("标题", chart_title, "")])
            slide = prs.slides.add_slide(blank); background(slide, colors["paper"])
            slide_header(
                slide, "DATA VIEW", chart_title_parts[0] + ("…" if len(chart_title_parts) > 1 else ""), len(prs.slides)
            )
            maximum = max(value for _, value in points) or 1
            available = 11.15; gap = available / max(1, len(points)); bar_width = min(0.72, gap * 0.58)
            for index, (label, value) in enumerate(points):
                height = max(0.05, 4.25 * value / maximum)
                left = 1.0 + index * gap + (gap - bar_width) / 2
                top = 6.15 - height
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(bar_width), Inches(height)
                )
                shape.fill.solid(); shape.fill.fore_color.rgb = colors["green"]; shape.line.fill.background()
                text_box(slide, f"{value:g}", left - 0.12, max(1.75, top - 0.32), bar_width + 0.24, 0.25,
                         size=9, bold=True, color=colors["moss"], align=PP_ALIGN.CENTER)
                text_box(slide, label[:14], left - 0.2, 6.28, bar_width + 0.4, 0.42,
                         size=9, color=colors["muted"], align=PP_ALIGN.CENTER)
            footer(slide)

        priority_labels = {"high": "高优先级", "medium": "中优先级", "low": "低优先级"}
        recommendation_records = []
        for index, item in enumerate(report.get("recommendations") or [], 1):
            action = _safe_text(item.get("action")); rationale = _safe_text(item.get("rationale"))
            recommendation_records.append((f"{index:02d}", f"{action} — {rationale}", priority_labels.get(item.get("priority"), "优先级待确认")))
        card_slides("ACTIONS", "建议行动", recommendation_records, "amber")

        limitations = report.get("limitations") or analysis.get("limitations") or []
        card_slides("LIMITATIONS", "限制与风险", [(f"{index:02d}", _safe_text(item), "需关注") for index, item in enumerate(limitations, 1)], "coral")
        next_steps = report.get("next_steps") or []
        if next_steps:
            card_slides("NEXT", "下一步", [(f"{index:02d}", _safe_text(item), "ACTION") for index, item in enumerate(next_steps, 1)], "green")

        target = report_dir / "report.pptx"
        prs.save(target)
        os.chmod(target, 0o600)
        return "ready"

    def get(self, report_id: str) -> dict[str, Any]:
        with self._connect() as conn:
            row = conn.execute("SELECT * FROM reports WHERE report_id = ?", (report_id,)).fetchone()
        if row is None:
            raise KeyError(report_id)
        return self._row(row)

    def _row(self, row: sqlite3.Row) -> dict[str, Any]:
        report_id, revision = row["report_id"], row["revision"]
        return {
            **dict(row),
            "internal_url": f"/reports/{report_id}",
            "technical_url": f"/reports/{report_id}/technical",
            "pdf_url": f"/reports/{report_id}/download/pdf" if row["pdf_status"] == "ready" else None,
            "pptx_url": f"/reports/{report_id}/download/pptx" if row["pptx_status"] == "ready" else None,
            "revision": revision,
        }

    def create_share(self, report_id: str, expires_at: str) -> dict[str, str]:
        self.get(report_id)
        try:
            expiry = datetime.fromisoformat(expires_at.replace("Z", "+00:00"))
        except ValueError as exc:
            raise ValueError("expires_at must be an RFC 3339 date-time") from exc
        if expiry.tzinfo is None:
            raise ValueError("expires_at must include a timezone")
        now = datetime.now(timezone.utc)
        if expiry <= now or (expiry - now).total_seconds() > 30 * 24 * 3600:
            raise ValueError("share expiry must be within the next 30 days")
        token = secrets.token_urlsafe(32)
        token_hash = hashlib.sha256(token.encode()).hexdigest()
        share_id = f"rps_{uuid4().hex}"
        with self._connect() as conn:
            conn.execute(
                "INSERT INTO report_shares VALUES (?, ?, ?, 'business', ?, NULL, ?)",
                (share_id, report_id, token_hash, expiry.isoformat(), _now()),
            )
        return {"share_id": share_id, "report_id": report_id, "token": token, "expires_at": expiry.isoformat()}

    def resolve_share(
        self, token: str, report_id: str | None = None, share_id: str | None = None
    ) -> dict[str, str]:
        token_hash = hashlib.sha256(token.encode()).hexdigest()
        with self._connect() as conn:
            row = conn.execute(
                "SELECT * FROM report_shares WHERE token_hash = ? AND revoked_at IS NULL",
                (token_hash,),
            ).fetchone()
        if row is None or (report_id is not None and row["report_id"] != report_id) or (
            share_id is not None and row["share_id"] != share_id
        ):
            raise KeyError("share")
        expiry = datetime.fromisoformat(row["expires_at"])
        if expiry <= datetime.now(timezone.utc):
            raise KeyError("share")
        return dict(row)

    def revoke_share(self, report_id: str, share_id: str) -> None:
        with self._connect() as conn:
            result = conn.execute(
                "UPDATE report_shares SET revoked_at=? WHERE report_id=? AND share_id=? AND revoked_at IS NULL",
                (_now(), report_id, share_id),
            )
        if result.rowcount != 1:
            raise KeyError(share_id)

    def file(self, report_id: str, name: str) -> Path:
        report = self.get(report_id)
        path = self.artifact_dir / report_id / f"v{report['revision']}" / name
        if not path.is_file():
            raise KeyError(name)
        return path

    def audit_download(self, report_id: str, fmt: str, actor: str) -> None:
        with self._connect() as conn:
            conn.execute("INSERT INTO report_download_audit VALUES (?, ?, ?, ?, ?)",
                         (f"rda_{uuid4().hex}", report_id, fmt, actor[:128], _now()))
